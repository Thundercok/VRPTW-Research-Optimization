# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_poster():
    # 1. Khởi tạo Presentation và thiết lập kích thước A0 Dọc (Vertical A0)
    # A0: 84.1 cm x 118.9 cm -> 33.11 inches x 46.81 inches
    prs = Presentation()
    prs.slide_width = Inches(33.11)
    prs.slide_height = Inches(46.81)

    # Sử dụng layout trống (blank slide)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 2. Định nghĩa bảng màu (Color Palette)
    COLOR_BG_DEEP = RGBColor(30, 15, 55)       # Tím đậm nền (#1E0F37)
    COLOR_HEADER_ACC = RGBColor(42, 22, 75)     # Tím sáng hơn ở header (#2A164B)
    COLOR_WHITE = RGBColor(255, 255, 255)       # Trắng nền nội dung (#FFFFFF)
    COLOR_BAR_YELLOW = RGBColor(254, 243, 214)  # Vàng nhạt thanh tiêu đề (#FEF3D6)
    COLOR_TEXT_NAVY = RGBColor(12, 30, 90)      # Xanh navy tiêu đề mục (#0C1E5A)
    COLOR_TEXT_DARK = RGBColor(51, 51, 51)      # Xám đậm chữ thường (#333333)
    COLOR_CYAN = RGBColor(0, 240, 255)          # Cyan công nghệ phát sáng (#00F0FF)
    COLOR_DECOR_DOT = RGBColor(78, 48, 143)     # Tím nhạt cho dot grid (#4E308F)
    COLOR_TEXT_MUTED = RGBColor(180, 170, 210)  # Tím nhạt chữ phụ đề (#B4AAD2)

    # 3. Vẽ nền tối toàn poster
    bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = COLOR_BG_DEEP
    bg_rect.line.fill.background() # Không viền

    # 4. Vẽ họa tiết trang trí góc / mạng lưới công nghệ ở nền (Header Decorative Grid & Lines)
    # Vẽ Dot Grid 5x4 ở phần trung tâm header
    grid_start_x = Inches(16.5)
    grid_start_y = Inches(1.5)
    for c in range(6):
        for r in range(4):
            cx = grid_start_x + Inches(c * 0.7)
            cy = grid_start_y + Inches(r * 0.7)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.18), Inches(0.18))
            dot.fill.solid()
            dot.fill.fore_color.rgb = COLOR_DECOR_DOT
            dot.line.fill.background()

    # Vẽ một vài đường nối công nghệ mảnh (Network node lines)
    def draw_tech_line(x1, y1, x2, y2):
        connector = slide.shapes.add_connector(1, x1, y1, x2, y2) # 1 = MSO_CONNECTOR.STRAIGHT
        connector.line.color.rgb = COLOR_DECOR_DOT
        connector.line.width = Pt(1.5)

    draw_tech_line(Inches(16.59), Inches(1.59), Inches(20.09), Inches(3.69))
    draw_tech_line(Inches(17.29), Inches(2.99), Inches(18.69), Inches(1.59))
    draw_tech_line(Inches(18.69), Inches(3.69), Inches(20.09), Inches(2.29))

    # Vẽ bốn vòng tròn trang trí ở mép trái phía dưới
    for i in range(4):
        cy = Inches(38.0) + Inches(i * 1.8)
        c_decor = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), cy, Inches(0.6), Inches(0.6))
        c_decor.fill.solid()
        c_decor.fill.fore_color.rgb = COLOR_DECOR_DOT
        c_decor.line.fill.background()

    # Họa tiết cyan ở viền dưới cùng (Representing circuitry)
    def draw_bottom_circuit():
        x_mid = prs.slide_width / 2
        y_bottom = prs.slide_height - Inches(1.0)
        
        # Thanh ngang chính giữa
        connector1 = slide.shapes.add_connector(1, x_mid - Inches(5.0), y_bottom, x_mid + Inches(5.0), y_bottom)
        connector1.line.color.rgb = COLOR_CYAN
        connector1.line.width = Pt(3.0)
        
        # Hai nhánh xéo trái và phải
        c_left = slide.shapes.add_connector(1, x_mid - Inches(5.0), y_bottom, x_mid - Inches(6.5), y_bottom - Inches(0.5))
        c_left.line.color.rgb = COLOR_CYAN
        c_left.line.width = Pt(2.0)
        
        c_right = slide.shapes.add_connector(1, x_mid + Inches(5.0), y_bottom, x_mid + Inches(6.5), y_bottom - Inches(0.5))
        c_right.line.color.rgb = COLOR_CYAN
        c_right.line.width = Pt(2.0)

    draw_bottom_circuit()

    # 5. Vẽ Khung Trắng lớn chứa nội dung chính (White Canvas Card)
    # Đặt lề trái/phải = 1.5 inch, lề trên = 10.2 inch để chừa chỗ cho Header tối màu, lề dưới = 1.5 inch
    canvas_left = Inches(1.5)
    canvas_top = Inches(10.2)
    canvas_width = prs.slide_width - (Inches(1.5) * 2) # 30.11 inches
    canvas_height = prs.slide_height - canvas_top - Inches(1.5) # 35.11 inches

    canvas_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, canvas_left, canvas_top, canvas_width, canvas_height)
    canvas_bg.fill.solid()
    canvas_bg.fill.fore_color.rgb = COLOR_WHITE
    canvas_bg.line.fill.background() # Không viền

    # 6. THIẾT KẾ PHẦN HEADER (TIÊU ĐỀ, LOGO & THÔNG TIN CHUNG)
    # 6.1. Dòng chữ phân ngành góc trái trên cùng
    fit_label_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(12.0), Inches(1.2))
    tf_fit = fit_label_box.text_frame
    tf_fit.word_wrap = True
    p_fit = tf_fit.paragraphs[0]
    p_fit.text = "NGHIÊN CỨU KHOA HỌC SINH VIÊN"
    p_fit.font.name = 'Arial'
    p_fit.font.size = Pt(36)
    p_fit.font.bold = True
    p_fit.font.color.rgb = COLOR_WHITE

    p_subfit = tf_fit.add_paragraph()
    p_subfit.text = "Khoa Công nghệ thông tin"
    p_subfit.font.name = 'Arial'
    p_subfit.font.size = Pt(24)
    p_subfit.font.italic = True
    p_subfit.font.color.rgb = COLOR_TEXT_MUTED
    p_subfit.space_before = Pt(6)

    # 6.2. Tiêu đề đề tài nghiên cứu chính
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(14.0), Inches(5.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    
    p_detai = tf_title.paragraphs[0]
    p_detai.text = "Đề tài"
    p_detai.font.name = 'Arial'
    p_detai.font.size = Pt(32)
    p_detai.font.bold = True
    p_detai.font.color.rgb = COLOR_WHITE
    p_detai.space_after = Pt(8)

    p_title_text = tf_title.add_paragraph()
    p_title_text.text = (
        "NGHIÊN CỨU VÀ TỐI ƯU HÓA BÀI TOÁN ĐỊNH TUYẾN PHƯƠNG TIỆN CÓ KHUNG THỜI GIAN (VRPTW) "
        "BẰNG THUẬT TOÁN TÌM KIẾM LÂN CẬN LỚN THÍCH ỨNG LAI HỌC TĂNG CƯỜNG SÂU (DDQN-ALNS)"
    )
    p_title_text.font.name = 'Arial'
    p_title_text.font.size = Pt(26)
    p_title_text.font.bold = True
    p_title_text.font.color.rgb = COLOR_WHITE
    p_title_text.space_after = Pt(12)

    # 6.3. Thông tin Tác giả (SVTH & GVHD) ở dưới tiêu đề
    p_author_info = tf_title.add_paragraph()
    p_author_info.text = "SVTH: Huỳnh Nhật Huy¹ (MSSV: 523c0012)\nGVHD: TS. Hồ Thị Linh*"
    p_author_info.font.name = 'Arial'
    p_author_info.font.size = Pt(18)
    p_author_info.font.bold = True
    p_author_info.font.color.rgb = COLOR_WHITE
    p_author_info.space_after = Pt(10)

    p_footnote = tf_title.add_paragraph()
    p_footnote.text = "¹ Sinh viên ngành Khoa học máy tính, Khoa Công nghệ thông tin\n* Bộ môn Khoa học máy tính, Khoa Công nghệ thông tin"
    p_footnote.font.name = 'Arial'
    p_footnote.font.size = Pt(14)
    p_footnote.font.italic = True
    p_footnote.font.color.rgb = COLOR_TEXT_MUTED

    # 6.4. Logo Đại học Tôn Đức Thắng góc phải trên cùng (Text-based stylized logo)
    logo_box = slide.shapes.add_textbox(Inches(21.0), Inches(0.8), Inches(10.5), Inches(1.5))
    tf_logo = logo_box.text_frame
    tf_logo.word_wrap = True
    
    p_uni = tf_logo.paragraphs[0]
    p_uni.text = "TRƯỜNG ĐẠI HỌC TÔN ĐỨC THẮNG"
    p_uni.alignment = PP_ALIGN.RIGHT
    p_uni.font.name = 'Arial'
    p_uni.font.size = Pt(24)
    p_uni.font.bold = True
    p_uni.font.color.rgb = COLOR_WHITE

    p_fac = tf_logo.add_paragraph()
    p_fac.text = "KHOA CÔNG NGHỆ THÔNG TIN"
    p_fac.alignment = PP_ALIGN.RIGHT
    p_fac.font.name = 'Arial'
    p_fac.font.size = Pt(20)
    p_fac.font.bold = True
    p_fac.font.color.rgb = COLOR_CYAN
    p_fac.space_before = Pt(4)

    # 6.5. Khung trắng góc phải đại diện cho QR Code / Ảnh
    photo_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(27.0), Inches(2.2), Inches(4.6), Inches(4.6))
    photo_box.fill.solid()
    photo_box.fill.fore_color.rgb = COLOR_WHITE
    photo_box.line.color.rgb = COLOR_DECOR_DOT
    photo_box.line.width = Pt(2.0)
    
    # Thêm văn bản hướng dẫn nhỏ bên trong khung ảnh
    tf_photo = photo_box.text_frame
    p_photo = tf_photo.paragraphs[0]
    p_photo.text = "[ QR CODE / ĐẠI DIỆN ĐỀ TÀI ]"
    p_photo.alignment = PP_ALIGN.CENTER
    p_photo.font.name = 'Arial'
    p_photo.font.size = Pt(14)
    p_photo.font.bold = True
    p_photo.font.color.rgb = COLOR_TEXT_DARK

    # 7. KHỞI TẠO CÁC PHÂN KHU NỘI DUNG (SECTION CREATOR HELPER)
    def add_section_header(title_text, left, top, width):
        # Vẽ thanh vàng nhạt
        header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.7))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = COLOR_BAR_YELLOW
        header_bg.line.fill.background()
        
        # Thêm text đè lên
        tb = slide.shapes.add_textbox(left, top, width, Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        # Tạo dấu ngoặc [] màu xanh lam sẫm
        run_bracket = p.add_run()
        run_bracket.text = "[] "
        run_bracket.font.name = 'Segoe UI'
        run_bracket.font.size = Pt(22)
        run_bracket.font.bold = True
        run_bracket.font.color.rgb = COLOR_TEXT_NAVY
        
        # Tạo text tiêu đề
        run_text = p.add_run()
        run_text.text = title_text
        run_text.font.name = 'Segoe UI'
        run_text.font.size = Pt(22)
        run_text.font.bold = True
        run_text.font.color.rgb = COLOR_TEXT_NAVY
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)

    def add_bullet(tf, title, body, font_size=15):
        p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = 1.2
        
        run_t = p.add_run()
        run_t.text = "• " + title + ": "
        run_t.font.name = 'Segoe UI'
        run_t.font.size = Pt(font_size)
        run_t.font.bold = True
        run_t.font.color.rgb = COLOR_TEXT_DARK
        
        run_b = p.add_run()
        run_b.text = body
        run_b.font.name = 'Segoe UI'
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = COLOR_TEXT_DARK

    # --- HÀNG 1: GIỚI THIỆU (CỘT TRÁI) & TẬP DỮ LIỆU (CỘT PHẢI) ---
    col1_left = Inches(2.0)
    col1_width = Inches(14.0)
    col2_left = Inches(17.11)
    col2_width = Inches(14.0)
    row1_top = Inches(10.8)
    row1_height = Inches(9.2)

    # 7.1. Khu vực GIỚI THIỆU
    add_section_header("GIỚI THIỆU", col1_left, row1_top, col1_width)
    intro_box = slide.shapes.add_textbox(col1_left, row1_top + Inches(0.8), col1_width, row1_height - Inches(0.8))
    tf_intro = intro_box.text_frame
    tf_intro.word_wrap = True
    tf_intro.paragraphs[0].text = "Bài toán định tuyến phương tiện có khung thời gian (VRPTW) đóng vai trò then chốt trong chuỗi cung ứng logistics đô thị hiện đại."
    tf_intro.paragraphs[0].font.name = 'Segoe UI'
    tf_intro.paragraphs[0].font.size = Pt(16)
    tf_intro.paragraphs[0].font.color.rgb = COLOR_TEXT_DARK
    tf_intro.paragraphs[0].space_after = Pt(8)

    add_bullet(tf_intro, "Phát biểu bài toán", "Xác định các tuyến đường tối ưu cho đội xe K xuất phát từ kho (depot) đi phục vụ khách hàng rải rác địa lý với nhu cầu q_i và khoảng khung thời gian phục vụ bắt buộc [E_i, L_i].")
    add_bullet(tf_intro, "Hàm mục tiêu MILP", "min f(x) = M · ∑_(k,j) x_0jk + ∑_(k,i,j) d_ij · x_ijk. Ưu tiên hàng đầu (hệ số phạt lớn M) là giảm quy mô đội xe vận hành (fleet size), và mục tiêu thứ cấp là giảm quãng đường đi.")
    add_bullet(tf_intro, "Hạn chế ALNS truyền thống", "Adaptive Large Neighborhood Search (ALNS) dựa trên các quy luật tính toán cứng nhắc để chọn toán tử, dễ mắc kẹt tại các cực trị địa phương và phụ thuộc nặng nề vào các tham số Simulated Annealing thiết lập thủ công.")
    add_bullet(tf_intro, "Giải pháp đề xuất", "Tích hợp Học tăng cường sâu (DRL) thông qua mô hình Double DQN và quy hoạch nguyên Set Partitioning để điều phối quá trình tìm kiếm thích ứng.")

    # 7.2. Khu vực TẬP DỮ LIỆU
    add_section_header("TẬP DỮ LIỆU & THIẾT LẬP THỰC NGHIỆM", col2_left, row1_top, col2_width)
    data_box = slide.shapes.add_textbox(col2_left, row1_top + Inches(0.8), col2_width, row1_height - Inches(0.8))
    tf_data = data_box.text_frame
    tf_data.word_wrap = True
    tf_data.paragraphs[0].text = "Nghiên cứu sử dụng bộ dữ liệu chuẩn Solomon và thiết lập phần cứng chuyên dụng để đánh giá hiệu năng thuật toán lai."
    tf_data.paragraphs[0].font.name = 'Segoe UI'
    tf_data.paragraphs[0].font.size = Pt(16)
    tf_data.paragraphs[0].font.color.rgb = COLOR_TEXT_DARK
    tf_data.paragraphs[0].space_after = Pt(8)

    add_bullet(tf_data, "Tập Solomon chính", "Đánh giá chi tiết trên các lớp RC1 (khung thời gian ngặt, khách hàng phân bố hỗn hợp giữa cụm và ngẫu nhiên) và RC2 (khung thời gian rộng, sức tải xe lớn). Quy mô thực thể gồm 100 khách hàng.")
    add_bullet(tf_data, "Ngẫu nhiên hóa miền (Domain Randomization)", "Huấn luyện chính sách DRL trên hệ thống mô phỏng dữ liệu tự động với quy mô từ 25 đến 100 khách hàng và các phân phối khung thời gian ngẫu nhiên để tăng khả năng tổng quát hóa.")
    add_bullet(tf_data, "Môi trường phần cứng", "CPU Apple M1 (Apple Silicon, 8 nhân), 16 GB bộ nhớ thống nhất (unified memory).")
    add_bullet(tf_data, "Môi trường phần mềm", "Python 3.10, PyTorch 2.1 (CPU-optimized), SciPy 1.11, Numba JIT (tối ưu hóa biên dịch trực tiếp các toán tử hình học và tìm kiếm).")

    # --- HÀNG 2: PHƯƠNG PHÁP THỰC HIỆN (BĂNG RỘNG TOÀN CHIỀU NGANG) ---
    row2_top = Inches(20.4)
    row2_width = prs.slide_width - (Inches(2.0) * 2) # 29.11 inches
    row2_height = Inches(12.2)

    add_section_header("PHƯƠNG PHÁP THỰC HIỆN (KIẾN TRÚC LAI DDQN-ALNS)", col1_left, row2_top, row2_width)
    
    # Chia nhỏ khu vực Phương pháp làm 3 cột thông tin để trông chuyên nghiệp hơn
    m_col_width = Inches(9.2)
    m_col_gap = Inches(0.7)
    
    # Cột Phương pháp 1: High-level & Low-level Control
    m_box_1 = slide.shapes.add_textbox(col1_left, row2_top + Inches(0.8), m_col_width, row2_height - Inches(0.8))
    tf_m1 = m_box_1.text_frame
    tf_m1.word_wrap = True
    tf_m1.paragraphs[0].text = "Hệ thống lai đề xuất sử dụng kiến trúc điều khiển hai tầng học sâu (Plateau & Operator Controllers) kết hợp tối ưu quy hoạch nguyên:"
    tf_m1.paragraphs[0].font.name = 'Segoe UI'
    tf_m1.paragraphs[0].font.size = Pt(15)
    tf_m1.paragraphs[0].font.bold = True
    tf_m1.paragraphs[0].font.color.rgb = COLOR_TEXT_NAVY
    tf_m1.paragraphs[0].space_after = Pt(10)
    
    add_bullet(tf_m1, "Plateau Controller (DRL cấp cao)", "Mạng Dueling Double DQN nhận vector trạng thái 12 chiều đặc trưng (tỷ lệ trì trệ, nhiệt độ, áp lực xe...) để chuyển đổi linh hoạt giữa 6 Search Mode (default, intensify, diversify, tw_rescue, pool_recombine, route_reduce).")
    add_bullet(tf_m1, "Operator Controller (DRL cấp thấp)", "Mạng Double DQN thứ hai nhận vector đặc trưng 15 chiều để chọn cặp toán tử tốt nhất từ 40 tổ hợp ALNS có sẵn. Kết hợp cơ chế UCB và Thompson Bandit để tăng cường tính khám phá.")
    add_bullet(tf_m1, "Hàm phần thưởng thế năng", "Thiết lập cơ chế thế năng thích ứng hướng xe: ưu tiên thưởng lớn cho hành động giảm số lượng phương tiện trước khi tối ưu khoảng cách.")

    # Cột Phương pháp 2: LAC & Set Partitioning
    m_box_2 = slide.shapes.add_textbox(col1_left + m_col_width + m_col_gap, row2_top + Inches(0.8), m_col_width, row2_height - Inches(0.8))
    tf_m2 = m_box_2.text_frame
    tf_m2.word_wrap = True
    tf_m2.paragraphs[0].text = "Sự kết hợp giữa học máy phân loại và quy hoạch toán học MILP giải quyết triệt để các ràng buộc ngặt:"
    tf_m2.paragraphs[0].font.name = 'Segoe UI'
    tf_m2.paragraphs[0].font.size = Pt(15)
    tf_m2.paragraphs[0].font.bold = True
    tf_m2.paragraphs[0].font.color.rgb = COLOR_TEXT_NAVY
    tf_m2.paragraphs[0].space_after = Pt(10)
    
    add_bullet(tf_m2, "Learned Acceptance Criterion (LAC)", "Mạng phân loại nhị phân 3 lớp thay thế Simulated Annealing. Được huấn luyện trực tuyến bằng thuật toán Weighted BCE và gán nhãn trễ (delayed labeling), giúp dự báo chính xác xác suất chấp nhận giải pháp ứng viên.")
    add_bullet(tf_m2, "Tái tổ hợp Set Partitioning", "Tất cả các tuyến đường tốt nhất được lưu trữ trong Route Pool. Định kỳ, mô hình giải bài toán phân hoạch tập hợp nguyên số (MILP) qua solver SciPy trong giới hạn 4 giây để lắp ghép giải pháp tối ưu toán học: min ∑ (P_veh + c_j) y_j thỏa mãn ∑ a_ij y_j = 1.")
    add_bullet(tf_m2, "Tìm kiếm cục bộ (Local Search)", "Áp dụng bộ lọc bán kính hạt mịn (granular filter) trên các toán tử 2-opt, Relocate, Swap, và Cross-Exchange nhằm tăng tốc.")

    # Cột Phương pháp 3: Operators & Route Reduction
    m_box_3 = slide.shapes.add_textbox(col1_left + (m_col_width + m_col_gap) * 2, row2_top + Inches(0.8), m_col_width, row2_height - Inches(0.8))
    tf_m3 = m_box_3.text_frame
    tf_m3.word_wrap = True
    tf_m3.paragraphs[0].text = "Các toán tử chuyên biệt được phát triển để xử lý ràng buộc khung thời gian và số xe:"
    tf_m3.paragraphs[0].font.name = 'Segoe UI'
    tf_m3.paragraphs[0].font.size = Pt(15)
    tf_m3.paragraphs[0].font.bold = True
    tf_m3.paragraphs[0].font.color.rgb = COLOR_TEXT_NAVY
    tf_m3.paragraphs[0].space_after = Pt(10)
    
    add_bullet(tf_m3, "8 Toán tử phá hủy", "Bao gồm phá hủy ngẫu nhiên, tệ nhất, tương đồng Shaw (địa lý, thời gian, tải trọng), một phần tuyến, khẩn cấp khung thời gian (tw_urgent), xóa tuyến toàn phần (route_eliminate), xóa tuyến phân tán, và Shaw liên tuyến.")
    add_bullet(tf_m3, "5 Toán tử sửa chữa", "Bao gồm chèn tham lan, regret-2, regret-3, tham lam ưu tiên thời gian, và đặc biệt là toán tử **FTS greedy** chèn khách hàng dựa trên việc cực tiểu hóa Composite Cost và bảo toàn Forward Time Slack (thời gian đệm còn lại của các khách hàng phía sau tuyến).")
    add_bullet(tf_m3, "Giảm xe chủ động", "Giải thuật loại bỏ tuyến nhỏ nhất và phân phối các khách hàng trên đó vào các tuyến khác khả thi, trực tiếp hạ số xe vận hành.")

    # --- HÀNG 3: KẾT QUẢ THỰC NGHIỆM (CỘT TRÁI) & KẾT LUẬN (CỘT PHẢI) ---
    col3_left = Inches(2.0)
    col3_width = Inches(18.5)
    col4_left = Inches(21.2)
    col4_width = Inches(9.91)
    row3_top = Inches(33.0)
    row3_height = Inches(11.0)

    # 7.3. Khu vực KẾT QUẢ THỰC NGHIỆM
    add_section_header("KẾT QUẢ THỰC NGHIỆM & PHÂN TÍCH", col3_left, row3_top, col3_width)
    
    # 7.3.1. Thêm bảng số liệu native trong PPTX
    # 11 hàng, 6 cột
    table_left = col3_left
    table_top = row3_top + Inches(0.8)
    table_width = col3_width
    table_height = Inches(4.5)
    
    table_shape = slide.shapes.add_table(11, 6, table_left, table_top, table_width, table_height)
    table = table_shape.table
    
    # Thiết lập độ rộng cột
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(4.2)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(3.3)
    table.columns[4].width = Inches(3.0)
    table.columns[5].width = Inches(3.0)
    
    headers = ["Lớp dữ liệu", "Thuật toán", "Số xe trung bình (NV)", "Quãng đường trung bình", "Độ lệch Gap%", "Thời gian chạy"]
    data = [
        ["RC1", "ALNS-Base", "12.575", "1327.91", "+1.909%", "19.2s"],
        ["RC1", "Hybrid-Fixed", "12.300", "1302.48", "-0.055%", "36.7s"],
        ["RC1", "Hybrid-Rule", "12.250", "1298.41", "-0.368%", "36.6s"],
        ["RC1", "Hybrid-DDQN (Đề xuất)", "12.350", "1298.54", "-0.358%", "40.3s"],
        ["RC1", "OR-Tools (CP-SAT)", "13.625", "1343.35", "+3.088%", "60.1s"],
        ["RC2", "ALNS-Base", "3.500", "1146.51", "+2.774%", "19.2s"],
        ["RC2", "Hybrid-Fixed", "3.400", "1131.42", "+1.425%", "36.7s"],
        ["RC2", "Hybrid-Rule", "3.400", "1128.16", "+1.130%", "36.6s"],
        ["RC2", "Hybrid-DDQN (Đề xuất)", "3.475", "1125.55", "+0.898%", "40.3s"],
        ["RC2", "OR-Tools (CP-SAT)", "6.250", "1034.02", "-7.350%*", "60.1s"]
    ]
    
    # Điền tiêu đề bảng
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_ACC
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Segoe UI'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        
    # Điền dữ liệu bảng
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            # Tô màu xen kẽ hàng
            if (row_idx + 1) % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(245, 245, 250)
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE
                
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Segoe UI'
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT_DARK
            # Highlight thuật toán đề xuất
            if "Hybrid-DDQN" in row_data[1]:
                p.font.bold = True
                p.font.color.rgb = RGBColor(12, 30, 150)
                if col_idx == 1:
                    cell.fill.fore_color.rgb = COLOR_BAR_YELLOW

    # Thêm ghi chú phân tích dưới bảng
    analysis_box = slide.shapes.add_textbox(col3_left, row3_top + Inches(5.4), col3_width, row3_height - Inches(5.4))
    tf_anal = analysis_box.text_frame
    tf_anal.word_wrap = True
    
    p_anal_title = tf_anal.paragraphs[0]
    p_anal_title.text = "Phân tích và bình luận chuyên sâu:"
    p_anal_title.font.name = 'Segoe UI'
    p_anal_title.font.size = Pt(15)
    p_anal_title.font.bold = True
    p_anal_title.font.color.rgb = COLOR_TEXT_NAVY
    p_anal_title.space_after = Pt(6)
    
    add_bullet(tf_anal, "Hiện tượng lạm phát xe (NV inflated) của OR-Tools", "Trên nhóm RC2, OR-Tools cho kết quả khoảng cách ngắn vượt trội (Gap% âm sâu -7.350%) nhưng bắt buộc sử dụng trung bình tới 6.25 xe (tăng 80% so với 3.475 xe của Hybrid-DDQN). Trong logistics đô thị thực tế, chi phí vận hành thêm một phương tiện cực lớn so với chi phí xăng xe chênh lệch, do đó giải pháp của Hybrid-DDQN hiệu quả kinh tế hơn rất nhiều.")
    add_bullet(tf_anal, "Đóng băng chính sách chuyển giao (Transfer-DR)", "Mô hình DRL sau khi huấn luyện Domain Randomization trên dữ liệu nhân tạo được đóng băng hoàn toàn trọng số. Khi giải trực tiếp 56 thực thể Solomon thu được Gap% cực kỳ ấn tượng 1.62% mà không cần cập nhật trực tuyến, chứng minh khả năng tổng quát hóa thực tiễn dạng plug-and-play.")
    add_bullet(tf_anal, "Cổng điều phối tương tác NAMI", "Hệ thống visualizer tích hợp backend FastAPI xử lý tối ưu hóa đa tiến trình và frontend trực quan hóa lộ trình các phương tiện trên bản đồ số, hiển thị biểu đồ hội tụ chi phí thời gian thực.")

    # 7.4. Khu vực KẾT LUẬN & ĐỀ NGHỊ
    add_section_header("KẾT LUẬN & HƯỚNG ĐỀ NGHỊ", col4_left, row3_top, col4_width)
    conclusion_box = slide.shapes.add_textbox(col4_left, row3_top + Inches(0.8), col4_width, row3_height - Inches(0.8))
    tf_concl = conclusion_box.text_frame
    tf_concl.word_wrap = True
    tf_concl.paragraphs[0].text = "Đề tài thiết lập thành công mô hình tối ưu lai DDQN-ALNS kết hợp thông minh học tăng cường sâu và quy hoạch toán học."
    tf_concl.paragraphs[0].font.name = 'Segoe UI'
    tf_concl.paragraphs[0].font.size = Pt(15)
    tf_concl.paragraphs[0].font.color.rgb = COLOR_TEXT_DARK
    tf_concl.paragraphs[0].space_after = Pt(8)

    add_bullet(tf_concl, "Đóng góp chính", "Phát triển bộ điều khiển Plateau & Operator bằng Dueling Double DQN; thay Simulated Annealing bằng cơ chế học LAC; phát triển toán tử sửa chữa bảo toàn Forward Time Slack (FTS) phối hợp Set Partitioning qua MILP.")
    add_bullet(tf_concl, "Hiệu năng thực tế", "Đạt Gap% chỉ 0.27% trên Solomon (ở cấu hình 1200 vòng lặp), duy trì quy mô đội xe tối thiểu tối ưu kinh tế vượt trội so với Google OR-Tools. Đóng gói hệ thống trực quan hóa điều phối NAMI.")
    add_bullet(tf_concl, "Hướng nghiên cứu tiếp theo", "Mở rộng bài toán hỗ trợ đội xe không đồng nhất (Heterogeneous Fleet VRP) với sức tải và định mức chi phí khác nhau; tích hợp các quy định bắt buộc về thời gian làm việc và nghỉ ngơi của tài xế; phát triển mô hình tái tối ưu hóa động (Dynamic Routing) thời gian thực.")

    # Thêm mục Tài liệu tham khảo rút gọn ở dưới cùng cột kết luận
    p_ref_title = tf_concl.add_paragraph()
    p_ref_title.text = "Tài liệu tham khảo:"
    p_ref_title.font.name = 'Segoe UI'
    p_ref_title.font.size = Pt(14)
    p_ref_title.font.bold = True
    p_ref_title.font.color.rgb = COLOR_TEXT_NAVY
    p_ref_title.space_before = Pt(12)
    p_ref_title.space_after = Pt(4)
    
    p_ref_text = tf_concl.add_paragraph()
    p_ref_text.text = (
        "1. Bi et al. (2022). A reinforcement learning-aided adaptive ALNS heuristic... IEEE Trans. Cybern., 52(9).\n"
        "2. Kool et al. (2021). Deep policy dynamic programming for vehicle routing. NeurIPS.\n"
        "3. Ropke & Pisinger (2006). An adaptive large neighborhood search heuristic... Transportation Science, 40(4)."
    )
    p_ref_text.font.name = 'Segoe UI'
    p_ref_text.font.size = Pt(11)
    p_ref_text.font.italic = True
    p_ref_text.font.color.rgb = COLOR_TEXT_DARK
    p_ref_text.line_spacing = 1.1

    # 8. Lưu tập tin presentation
    output_filename = "docs/Bao_Cao_Nghien_Cuu_VRPTW_Poster.pptx"
    prs.save(output_filename)
    print(f"Successfully generated scientific poster at: {os.path.abspath(output_filename)}")

if __name__ == '__main__':
    create_poster()
