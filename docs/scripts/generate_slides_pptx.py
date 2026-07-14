# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slides():
    prs = Presentation()
    
    # Set to 16:9 widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ═══════════════════════════════════════════════════════════════════
    # STYLE SYSTEM & THEME
    # ═══════════════════════════════════════════════════════════════════
    COLOR_BG = RGBColor(15, 23, 42)          # Deep slate background (#0f172a)
    COLOR_CARD = RGBColor(30, 41, 59)        # Dark gray cards (#1e293b)
    COLOR_WHITE = RGBColor(255, 255, 255)    # White
    COLOR_LIGHT_GRAY = RGBColor(226, 232, 240) # Text body (#e2e8f0)
    COLOR_DIM_GRAY = RGBColor(148, 163, 184) # Muted text (#94a3b8)
    COLOR_CYAN = RGBColor(0, 180, 216)       # Primary Cyan accent
    COLOR_GOLD = RGBColor(245, 158, 11)      # Secondary Gold accent
    COLOR_SUCCESS = RGBColor(16, 185, 129)   # Green for good results
    COLOR_DANGER = RGBColor(239, 68, 68)     # Red for warnings/fails
    
    FONT_TITLE = "Segoe UI"
    FONT_BODY = "Segoe UI"

    # Helpers
    def apply_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text=None):
        # Category Tag
        if category_text:
            cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
            tf = cat_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = category_text.upper()
            p.font.name = FONT_TITLE
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
            
        # Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    def add_card(slide, left, top, width, height, title=None, border_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
            
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = FONT_TITLE
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            
        return card

    def add_bullet_list(slide, left, top, width, height, bullets, font_size=14):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        for idx, b in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            # Handle indent
            if b.startswith("  - ") or b.startswith("    * "):
                p.level = 1
                p.text = b.split("- ", 1)[-1] if "- " in b else b.split("* ", 1)[-1]
            elif b.startswith("    - "):
                p.level = 2
                p.text = b.split("- ", 1)[-1]
            else:
                p.level = 0
                # strip out initial number or dot if present
                p.text = b
                
            p.font.name = FONT_BODY
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLOR_LIGHT_GRAY
            p.space_after = Pt(10)
            
        return tb

    def add_stats_box(slide, left, top, width, height, value, label):
        # Outer border card
        add_card(slide, left, top, width, height, border_color=COLOR_CYAN)
        
        # Value
        v_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.6))
        tf = v_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        
        # Label
        l_box = slide.shapes.add_textbox(left, top + Inches(0.8), width, Inches(0.4))
        tf_l = l_box.text_frame
        tf_l.word_wrap = True
        p_l = tf_l.paragraphs[0]
        p_l.text = label.upper()
        p_l.alignment = PP_ALIGN.CENTER
        p_l.font.name = FONT_BODY
        p_l.font.size = Pt(9)
        p_l.font.bold = True
        p_l.font.color.rgb = COLOR_DIM_GRAY

    def set_slide_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    def add_route_line(slide, x1, y1, x2, y2, color, width=2):
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = color
        conn.line.width = Pt(width)
        return conn

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE SLIDE
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    
    # Large Decorative Circle for title slide
    decor_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(4.5), Inches(4.5))
    decor_circle.fill.solid()
    decor_circle.fill.fore_color.rgb = RGBColor(18, 30, 60)
    decor_circle.line.fill.background()
    
    # Subtitle Category
    conf_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.4))
    tf = conf_box.text_frame
    p = tf.paragraphs[0]
    p.text = "STUDENT RESEARCH CONFERENCE (NCKH) 2025 - 2026"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10.5), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "A Hybrid Deep Reinforcement Learning and Adaptive Large Neighborhood Search Metaheuristic for the Vehicle Routing Problem with Time Windows"
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(12)
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(10.5), Inches(0.8))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Combining Neural Decision Making with Heuristics to Resolve Scale-Aware Routing Limits"
    p.font.name = FONT_BODY
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_GOLD
    
    # Meta layout
    meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2))
    tf = meta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Presenters: Thi-Linh Ho & Ha-Vy Duy Nguyen\nAdvisor: Student Research Initiative Committee\nFaculty of Information Technology, Ton Duc Thang University"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_LIGHT_GRAY
    p.font.bold = True
    p.space_after = Pt(6)
    
    set_slide_notes(slide, (
        "Kính thưa quý thầy cô trong Hội đồng khoa học và các bạn sinh viên.\n\n"
        "Hôm nay, em xin phép đại diện nhóm nghiên cứu trình bày đề tài: "
        "'Thuật toán lai học tăng cường sâu và tìm kiếm lân cận lớn thích ứng cho bài toán "
        "tối ưu hóa định tuyến xe có ràng buộc khung thời gian (VRPTW)'.\n\n"
        "Nghiên cứu của tụi em tập trung vào việc giải quyết bài toán cốt lõi trong ngành logistics "
        "bằng cách kết hợp sức mạnh của học máy (Deep RL) và toán học tối ưu hóa truyền thống (Metaheuristics) "
        "để vượt qua các rào cản về hiệu năng và số lượng xe của các giải pháp hiện tại."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: AGENDA
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Presentation Agenda", "Structure")
    
    agenda_items = [
        ("01. Introduction", "VRPTW parameters, constraints, and the Lexicographical objective model."),
        ("02. Motivation", "Addressing Fleet Inflation & the necessity of real-time telemetry."),
        ("03. Proposed Framework", "Combining Hierarchical MDP (DDQN) with mathematical search (ALNS)."),
        ("04. Proactive Seeding", "Clarke-Wright bridge routing & Set Partitioning MILP solver."),
        ("05. GNN Heatmaps", "Optimizing detour calculations dynamically with neural edge predictions."),
        ("06. Post-Processing", "Ejection chains, buffered route elimination, and TD polishing cascade."),
        ("07. System Architecture", "Monorepo integration of Vite + React UI, FastAPI, and C++ solver bindings."),
        ("08. Experimental Results", "Verification under strict cold-starts, OR-Tools comparison, and scaling limits.")
    ]
    
    card_w = Inches(2.7)
    card_h = Inches(2.1)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    
    for idx, (title, desc) in enumerate(agenda_items):
        row = idx // 4
        col = idx % 4
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        
        # Add Card
        add_card(slide, x, y, card_w, card_h)
        
        # Text
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Title
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.space_after = Pt(10)
        
        # Description
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_LIGHT_GRAY
        
    set_slide_notes(slide, (
        "Nội dung bài báo cáo hôm nay của em sẽ đi qua 8 phần chính:\n"
        "từ việc xác định bài toán, phân tích khoảng trống nghiên cứu, "
        "đề xuất mô hình lai DDQN-ALNS tích hợp GNN, thiết kế hệ thống NAMI Dispatcher thực tế, "
        "các kết quả thực nghiệm kiểm định nghiêm ngặt, so sánh trực diện với Google OR-Tools "
        "và cuối cùng là kết luận."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: PROBLEM DEFINITION
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Vehicle Routing Problem with Time Windows (VRPTW)", "Problem Modeling")
    
    # Left Card - Math
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Mathematical Foundations")
    bullets_s3 = [
        "Objective Context: Serve a set of customer demands q_i at locations N with homogeneous fleet capacity Q.",
        "Hard Operational Constraints:",
        "  - Capacity limit: Sum of demands along any route <= Q.",
        "  - Time Windows: Arrival at customer i must occur within [e_i, l_i]. Early arrival induces waiting; late arrival is strictly infeasible.",
        "Lexicographical Objective Hierarchy:",
        "    * Primary: Minimize Fleet Size (NV)",
        "    * Secondary: Minimize Travel Distance (TD)"
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s3, font_size=13)
    
    # Add inline mathematical formula representation
    f_box = slide.shapes.add_textbox(Inches(1.1), Inches(5.6), Inches(5.4), Inches(0.8))
    tf = f_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "lex min ( NV ,  TD )"
    p.font.name = FONT_BODY
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_CYAN
    
    # Right Card - Graph Diorama (PowerPoint Shapes!)
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Routing Diorama Visualization")
    
    # Draw Depot
    depot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.4), Inches(3.8), Inches(0.8), Inches(0.8))
    depot.fill.solid()
    depot.fill.fore_color.rgb = COLOR_GOLD
    depot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(8.4), Inches(4.6), Inches(2.8), Inches(0.5))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "DEPOT [0, 1000]"
    p.font.name = FONT_TITLE
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Draw Nodes
    nodes = [
        ("c1 [10, 50]", Inches(8.0), Inches(2.2), COLOR_CYAN),
        ("c2 [40, 90]", Inches(7.8), Inches(5.2), COLOR_CYAN),
        ("c3 [20, 60]", Inches(11.2), Inches(2.4), COLOR_CYAN),
        ("c4 [70, 120]", Inches(11.0), Inches(5.0), COLOR_CYAN),
    ]
    for name, nx, ny, color in nodes:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, nx, ny, Inches(0.4), Inches(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        
        lbl = slide.shapes.add_textbox(nx - Inches(0.8), ny + Inches(0.4), Inches(2.0), Inches(0.4))
        lbl.text_frame.word_wrap = True
        p = lbl.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_LIGHT_GRAY
        
    # Route 1: Depot (9.8, 4.2) -> c1 (8.2, 2.4) -> c2 (8.0, 5.4) -> Depot
    add_route_line(slide, Inches(9.8), Inches(4.2), Inches(8.2), Inches(2.4), COLOR_CYAN, 2.5)
    add_route_line(slide, Inches(8.2), Inches(2.4), Inches(8.0), Inches(5.4), COLOR_CYAN, 2.5)
    add_route_line(slide, Inches(8.0), Inches(5.4), Inches(9.8), Inches(4.2), COLOR_CYAN, 2.5)
    
    # Route 2: Depot (9.8, 4.2) -> c3 (11.4, 2.6) -> c4 (11.2, 5.2) -> Depot
    add_route_line(slide, Inches(9.8), Inches(4.2), Inches(11.4), Inches(2.6), COLOR_GOLD, 2.5)
    add_route_line(slide, Inches(11.4), Inches(2.6), Inches(11.2), Inches(5.2), COLOR_GOLD, 2.5)
    add_route_line(slide, Inches(11.2), Inches(5.2), Inches(9.8), Inches(4.2), COLOR_GOLD, 2.5)
        
    set_slide_notes(slide, (
        "VRPTW là một bài toán NP-Khó điển hình trong vận tải. "
        "Mục tiêu của chúng ta là điều phối một đội xe đồng nhất từ kho trung tâm đi phục vụ khách hàng sao cho: "
        "tổng số lượng xe (fleet size) được tối thiểu hóa trước tiên, sau đó mới đến tối thiểu hóa tổng quãng đường di chuyển (travel distance).\n\n"
        "Đây là mục tiêu phân cấp (lexicographical) rất sát với thực tế, vì chi phí cố định để vận hành thêm một chiếc xe lớn "
        "hơn rất nhiều so với chi phí nhiên liệu phát sinh thêm trên quãng đường."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4: RESEARCH MOTIVATION
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Research Motivation & Operational Gaps", "Motivation")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Fleet Economics & Logistics Objectives")
    bullets_s4 = [
        "Fleet Costs Dominate: Adding a single vehicle incurs fixed labor, capital overhead, and licensing costs equivalent to 10x to 50x the variable fuel cost of marginal travel distance changes.",
        "Lexicographical Entrapment: Standard algorithms trade vehicle count (increasing fleet size) to yield a slightly lower travel distance, causing 'fleet inflation'.",
        "Actionable Tooling Needs: Dispatch operators require interactive telemetry, low-latency Gantt diagrams and map rendering to manage live deliveries instead of plain solver console text."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s4, font_size=14)
    
    # Right Stats & Warn
    add_stats_box(slide, Inches(7.1), Inches(1.6), Inches(2.5), Inches(1.5), "$10,000+", "Fixed Cost / Month")
    add_stats_box(slide, Inches(10.0), Inches(1.6), Inches(2.5), Inches(1.5), "10x - 50x", "Savings Factor / Vehicle")
    
    # Warning Card
    warn_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(3.4), Inches(5.4), Inches(3.2))
    warn_card.fill.solid()
    warn_card.fill.fore_color.rgb = COLOR_CARD
    warn_card.line.color.rgb = COLOR_DANGER
    warn_card.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(Inches(7.4), Inches(3.6), Inches(4.8), Inches(2.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "⚠ The \"Fleet Inflation\" Trap"
    p.font.name = FONT_TITLE
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_DANGER
    p.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = "Traditional metaheuristics optimize scalar values: Cost = alpha * NV + beta * TD.\n\nIf alpha is set too low relative to time bounds, the algorithm sacrifices vehicle bounds to optimize travel coordinates. Our hybrid model separates these objectives strictly."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_LIGHT_GRAY
    
    set_slide_notes(slide, (
        "Động lực lớn nhất của nghiên cứu này là giải quyết hiện tượng 'lạm phát đội xe' (fleet inflation). "
        "Nhiều thuật toán hiện tại khi chạy thực tế thường chấp nhận tăng thêm xe để đổi lấy tổng quãng đường ngắn hơn một chút, "
        "điều này đi ngược lại tối ưu hóa chi phí vận hành doanh nghiệp.\n\n"
        "Đồng thời, tụi em muốn tạo ra một hệ thống không chỉ giải toán trên giấy mà phải tích hợp trực quan hóa thời gian thực (telemetry) "
        "để điều phối viên có thể giám sát và can thiệp khi cần thiết."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 5: RELATED WORK & RESEARCH GAPS
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Related Work & Research Gaps", "Literature Review")
    
    # Left Card with Table
    add_card(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0), "Comparative Analysis of VRP Methodologies")
    
    # Add Table
    rows = 4
    cols = 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(1.1), Inches(2.3), Inches(11.1), Inches(3.0))
    table = table_shape.table
    
    # Headers
    headers = ["Methodology", "Constraint Feasibility", "Scaling Capabilities", "Search Adaptability", "Objective Priority"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    data = [
        ["Classical Heuristics (ALNS, TS)", "100% Guaranteed", "High (1000+ nodes)", "Low (Static Rules)", "Scalar Weighted Cost"],
        ["End-to-End Deep Learning", "Fails on tight bounds", "Poor (< 100 nodes)", "High (Neural)", "Distance Minimization"],
        ["Proposed Hybrid Model", "100% Guaranteed", "High (400+ nodes)", "High (Hierarchical MDP)", "Lexicographical Primary"]
    ]
    
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = FONT_BODY
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_LIGHT_GRAY
            p.alignment = PP_ALIGN.CENTER
            
            # Highlight our row
            if r_idx == 2:
                p.font.bold = True
                if c_idx == 0:
                    p.font.color.rgb = COLOR_GOLD
                elif c_idx in [1, 2, 3]:
                    p.font.color.rgb = COLOR_SUCCESS
                else:
                    p.font.color.rgb = COLOR_CYAN

    # Gap Text Box
    gap_box = slide.shapes.add_textbox(Inches(1.1), Inches(5.6), Inches(11.1), Inches(0.8))
    tf = gap_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Core Research Gap: How can we coordinate a mathematically robust search that maintains 100% operational feasibility while utilizing neural network policies to govern destroy-repair transitions?"
    p.font.name = FONT_BODY
    p.font.italic = True
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_GOLD
    
    set_slide_notes(slide, (
        "Nhìn lại các nghiên cứu trước đây: Các phương pháp heuristic truyền thống như ALNS rất mạnh về đảm bảo ràng buộc, "
        "nhưng cơ chế chọn toán tử bằng vòng quay roulette lại quá đơn giản và mang tính myopic (thiển cận) — không học được trạng thái tìm kiếm hiện tại.\n\n"
        "Ngược lại, các mô hình học máy cuối-đến-cuối (End-to-End) như Pointer Networks hay GNN sinh lộ trình trực tiếp lại rất dễ tạo ra "
        "các tuyến đường không khả thi khi gặp ràng buộc thời gian ngặt nghèo của VRPTW. Vì vậy, nhóm đề xuất giải pháp lai: "
        "Giữ nguyên khung thuật toán ALNS để đảm bảo 100% tính khả thi toán học, nhưng dùng mạng thần kinh sâu DDQN "
        "làm bộ não điều khiển việc chọn toán tử và chấp nhận lời giải."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 6: METHODOLOGY - HIERARCHICAL MDP
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Hierarchical MDP Solver Core", "Methodology")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Dual-Level Decision Process")
    bullets_s6 = [
        "Macro-Level (Plateau Controller):",
        "  - Monitors solver convergence. Triggers if search stagnates for 60 iterations.",
        "  - Observes state s_t_c in R_12; selects search regime goal mode m_t in {Intensify, Diversify, TW-Rescue, Recombine, Reduce}.",
        "Micro-Level (Operator Controller):",
        "  - Triggers at every search step. Observes state s_t_o in R_15.",
        "  - Selects destroy-repair operator pair a_t in {0, ..., 54}.",
        "Thompson Bandit Integration: Decays and balances operator likelihoods dynamically based on real-time success feedback."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s6, font_size=13)
    
    # Right Card - Flowchart
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Hierarchical Information Flow")
    
    # Draw Macro Box
    add_card(slide, Inches(7.5), Inches(2.2), Inches(4.6), Inches(0.9), border_color=COLOR_GOLD)
    tb_mac = slide.shapes.add_textbox(Inches(7.5), Inches(2.3), Inches(4.6), Inches(0.7))
    tf = tb_mac.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Plateau Controller (Macro DDQN)"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Monitors Stagnation | State s_t_c (12D)"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_GOLD
    
    # Draw Connector 1
    add_route_line(slide, Inches(9.8), Inches(3.1), Inches(9.8), Inches(3.7), COLOR_GOLD, 2)
    
    # Draw Micro Box
    add_card(slide, Inches(7.5), Inches(3.7), Inches(4.6), Inches(0.9), border_color=COLOR_CYAN)
    tb_mic = slide.shapes.add_textbox(Inches(7.5), Inches(3.8), Inches(4.6), Inches(0.7))
    tf = tb_mic.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Operator Controller (Micro DDQN)"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Selects Heuristics | State s_t_o (15D)"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_CYAN
    
    # Draw Connector 2
    add_route_line(slide, Inches(9.8), Inches(4.6), Inches(9.8), Inches(5.2), COLOR_CYAN, 2)
    
    # Draw Env Box
    add_card(slide, Inches(7.9), Inches(5.2), Inches(3.8), Inches(0.8), border_color=COLOR_LIGHT_GRAY)
    tb_env = slide.shapes.add_textbox(Inches(7.9), Inches(5.4), Inches(3.8), Inches(0.5))
    tf = tb_env.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "ALNS Local Search Loop"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Labels
    lbl1 = slide.shapes.add_textbox(Inches(10.0), Inches(3.2), Inches(2.0), Inches(0.4))
    lbl1.text_frame.paragraphs[0].text = "m_t (Mode)"
    lbl1.text_frame.paragraphs[0].font.size = Pt(11)
    lbl1.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD
    
    lbl2 = slide.shapes.add_textbox(Inches(10.0), Inches(4.7), Inches(2.0), Inches(0.4))
    lbl2.text_frame.paragraphs[0].text = "a_t (Operators)"
    lbl2.text_frame.paragraphs[0].font.size = Pt(11)
    lbl2.text_frame.paragraphs[0].font.color.rgb = COLOR_CYAN

    set_slide_notes(slide, (
        "Tụi em thiết kế bài toán dưới dạng Tiến trình Quyết định Markov Phân cấp (Hierarchical MDP).\n\n"
        "Bộ điều khiển cấp cao Plateau Controller sẽ theo dõi tiến trình hội tụ. "
        "Nếu quá trình tìm kiếm bị kẹt (plateau) quá 60 vòng lặp, nó sẽ can thiệp bằng cách chuyển đổi chế độ tìm kiếm: "
        "ví dụ chuyển sang 'TW-Rescue' để gỡ các nút thắt thời gian, hoặc 'Reduce' để ép giảm xe.\n\n"
        "Dưới quyền của nó, bộ điều khiển cấp thấp Operator Controller sẽ căn cứ vào chế độ đang chạy "
        "và trạng thái chi tiết của các tuyến đường để chọn ra cặp toán tử phá hủy và sửa chữa tốt nhất trong số 55 cặp có sẵn."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 7: LEARNED ACCEPTANCE & RECOMBINATION
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Set Partitioning Route Recombination", "Methodology")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Neural Acceptance & Hybrid Optimization")
    bullets_s7 = [
        "Learned Acceptance Criterion (LAC):",
        "  - Replaces traditional, static Simulated Annealing (SA) criteria.",
        "  - Classifier network trained via online hindsight labels: predicts whether accepting a degraded search move is likely to lead to downstream improvements within horizon H = 80.",
        "Set Partitioning Route Recombination:",
        "  - Extracted routes from a dual-retention pool are solved periodically using MILP formulations.",
        "  - Goal: Choose the subset of elite routes that covers all customers exactly once while minimizing vehicle count and distance.",
        "  - Solved with SciPy MILP under a strict 4.0s limit, scaling vehicle penalty P_vehicle to force fleet drops."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s7, font_size=12)
    
    # Right Card
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Dual-Retention Route Pool")
    
    # Box 1
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.4), Inches(4.6), Inches(1.2))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(40, 25, 60)
    box1.line.color.rgb = COLOR_GOLD
    box1.line.width = Pt(1.5)
    tb1 = slide.shapes.add_textbox(Inches(7.7), Inches(2.5), Inches(4.2), Inches(1.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "1. Elite Route Memory"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p2 = tf1.add_paragraph()
    p2.text = "Caches disjoint routes from the top-scoring solutions found throughout the run, preserving high-quality sub-routes."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_LIGHT_GRAY
    
    # Box 2
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3.9), Inches(4.6), Inches(1.2))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(20, 45, 65)
    box2.line.color.rgb = COLOR_CYAN
    box2.line.width = Pt(1.5)
    tb2 = slide.shapes.add_textbox(Inches(7.7), Inches(4.0), Inches(4.2), Inches(1.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "2. Diversity Route Cache"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p2 = tf2.add_paragraph()
    p2.text = "Tracks distinct routes that utilize alternative links to guarantee routing variety during optimization."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_LIGHT_GRAY
    
    # Footnote Box
    add_card(slide, Inches(7.5), Inches(5.4), Inches(4.6), Inches(0.8), border_color=COLOR_DIM_GRAY)
    tb3 = slide.shapes.add_textbox(Inches(7.7), Inches(5.5), Inches(4.2), Inches(0.6))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "★ Scalability Priority: Limits candidate pool size to 500 routes to ensure MILP is solved in real-time."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    set_slide_notes(slide, (
        "Một cải tiến quan trọng khác là Learned Acceptance Criterion (LAC).\n\n"
        "Thay vì chấp nhận lời giải tệ hơn bằng công thức Simulated Annealing vật lý truyền thống, "
        "chúng em dùng một mạng thần kinh phân loại nhị phân siêu nhẹ. Mạng này học trực tuyến qua cơ chế hindsight (nhìn lại quá khứ) "
        "để đoán xem: việc chấp nhận một bước đi lùi hiện tại có giúp thuật toán tìm ra đột phá trong vòng 80 bước tiếp theo hay không.\n\n"
        "Bên cạnh đó, các tuyến đường tốt đi qua các vòng lặp được lưu trữ vào Route Pool. Định kỳ, mô hình giải bài toán phân hoạch tập hợp (Set Partitioning) "
        "bằng MILP dưới giới hạn 4 giây để tìm ra cách lắp ghép các tuyến đường đơn lẻ thành một lời giải hoàn chỉnh tối ưu nhất về số xe."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 8: PROACTIVE POOL SEEDING & COLUMN GENERATION
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Proactive Pool Seeding & Route Generation", "Methodology")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Overcoming Feasibility Entrapment")
    bullets_s8 = [
        "The Challenge: Traditional local search generates routes reactively. Difficult instances (e.g. Solomon RC101) require highly complex 'bridge routes' to reduce vehicle count, which small mutations fail to build.",
        "Proactive Seeding Mechanisms:",
        "  1. Randomized Clarke-Wright Savings: Constructs temporal bridge routes between spatially distant but highly compatible nodes.",
        "  2. NV-Targeted Construction: Forces greedy routing variants directly aiming for NV - 1 vehicle counts, mapping remaining stops via high-capacity models.",
        "  3. Regret-3 Large Destroy (40-65%): Removes up to 65% of the graph, forcing extensive reorganizations to pack routes."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s8, font_size=12.5)
    
    # Right Card - Flowchart
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Column Construction Pipelines")
    
    # 3 Sources
    add_card(slide, Inches(7.4), Inches(2.3), Inches(1.4), Inches(0.8))
    tb_s1 = slide.shapes.add_textbox(Inches(7.4), Inches(2.4), Inches(1.4), Inches(0.6))
    tb_s1.text_frame.word_wrap = True
    tb_s1.text_frame.paragraphs[0].text = "Clarke-Wright\nSavings"
    tb_s1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_s1.text_frame.paragraphs[0].font.size = Pt(11)
    tb_s1.text_frame.paragraphs[0].font.color.rgb = COLOR_CYAN
    tb_s1.text_frame.paragraphs[0].font.bold = True
    
    add_card(slide, Inches(9.1), Inches(2.3), Inches(1.4), Inches(0.8))
    tb_s2 = slide.shapes.add_textbox(Inches(9.1), Inches(2.4), Inches(1.4), Inches(0.6))
    tb_s2.text_frame.word_wrap = True
    tb_s2.text_frame.paragraphs[0].text = "NV-Targeted\nConstruction"
    tb_s2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_s2.text_frame.paragraphs[0].font.size = Pt(11)
    tb_s2.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD
    tb_s2.text_frame.paragraphs[0].font.bold = True
    
    add_card(slide, Inches(10.8), Inches(2.3), Inches(1.4), Inches(0.8))
    tb_s3 = slide.shapes.add_textbox(Inches(10.8), Inches(2.4), Inches(1.4), Inches(0.6))
    tb_s3.text_frame.word_wrap = True
    tb_s3.text_frame.paragraphs[0].text = "Large Destroy\n(40% - 65%)"
    tb_s3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_s3.text_frame.paragraphs[0].font.size = Pt(11)
    tb_s3.text_frame.paragraphs[0].font.color.rgb = COLOR_DANGER
    tb_s3.text_frame.paragraphs[0].font.bold = True
    
    # Connecting Lines to Route Pool
    add_route_line(slide, Inches(8.1), Inches(3.1), Inches(9.8), Inches(3.7), COLOR_DIM_GRAY, 1.5)
    add_route_line(slide, Inches(9.8), Inches(3.1), Inches(9.8), Inches(3.7), COLOR_DIM_GRAY, 1.5)
    add_route_line(slide, Inches(11.5), Inches(3.1), Inches(9.8), Inches(3.7), COLOR_DIM_GRAY, 1.5)
    
    # Route Pool
    add_card(slide, Inches(7.7), Inches(3.7), Inches(4.2), Inches(0.9), border_color=COLOR_GOLD)
    tb_pool = slide.shapes.add_textbox(Inches(7.7), Inches(3.8), Inches(4.2), Inches(0.7))
    tf = tb_pool.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Candidate Route Pool"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Dual Elite & Diversity Splits"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_GOLD
    
    # Line to MILP
    add_route_line(slide, Inches(9.8), Inches(4.6), Inches(9.8), Inches(5.1), COLOR_GOLD, 2)
    
    # MILP Solver
    add_card(slide, Inches(7.7), Inches(5.1), Inches(4.2), Inches(0.9), border_color=COLOR_CYAN)
    tb_milp = slide.shapes.add_textbox(Inches(7.7), Inches(5.2), Inches(4.2), Inches(0.7))
    tf = tb_milp.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Set Partitioning (MILP)"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "SciPy Optimizer (< 4s Limit)"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_CYAN
    
    # Line to final result
    add_route_line(slide, Inches(9.8), Inches(6.0), Inches(9.8), Inches(6.3), COLOR_CYAN, 2)
    
    # Result text
    lbl = slide.shapes.add_textbox(Inches(7.7), Inches(6.3), Inches(4.2), Inches(0.3))
    lbl.text_frame.paragraphs[0].text = "MINIMIZED VEHICLE (NV) PLAN"
    lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    lbl.text_frame.paragraphs[0].font.size = Pt(12)
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.color.rgb = COLOR_SUCCESS
    
    set_slide_notes(slide, (
        "Khi tụi em thử nghiệm trên các thực thể cực khó như Solomon RC101, "
        "thuật toán thường bị kẹt ở mức 15 xe trong khi lời giải tốt nhất thế giới (BKS) là 14 xe.\n\n"
        "Qua phân tích, chúng em phát hiện ra lỗi Feasibility Entrapment: "
        "thuật toán ALNS thông thường chỉ cải tiến cục bộ nên không bao giờ tạo ra được các tuyến đường có cấu trúc "
        "hoàn toàn mới để gom nhóm khách hàng khác đi. Để phá vỡ trần xe này, nhóm đã bổ sung cơ chế gieo hạt chủ động "
        "vào Route Pool: sử dụng thuật toán Clarke-Wright Savings cải tiến để tạo 'cầu nối thời gian', "
        "khởi tạo trực tiếp hướng mục tiêu số xe NV-1, và thực hiện sinh cột thuật toán từ lời giải tốt nhất.\n\n"
        "Nhờ vậy, Route Pool luôn giàu các tuyến tiềm năng dài 8-10 khách hàng để bộ giải MILP tìm ra lời giải 14 xe."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 9: GNN EDGE HEATMAP GUIDANCE
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Graph Neural Network Guided Search", "Methodology")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Dynamic Edge Optimization")
    bullets_s9 = [
        "Neural Heatmap Predictor: A lightweight Graph Neural Network predicts edge optimality probabilities P_ij in [0, 1] exactly once during initialization.",
        "GNN-Biased Heuristics: Re-scales detour insertion costs dynamically during regret/greedy heuristics:",
        "  - Formula: detour' = detour * (1.0 - gamma * P_iu) * (1.0 - gamma * P_uj)",
        "  - Where gamma in [0.1, 0.6] is scaled dynamically based on active search mode.",
        "Search Pruning Schedule: Linear threshold decay (theta = 0.05 -> 0.003) eliminates unlikely edge transitions, pruning the search space to accelerate large instances."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s9, font_size=12.5)
    
    # Right Card - Heatmap visualization (PowerPoint Shapes!)
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "GNN Probability Heatmap")
    
    # Nodes on Graph
    gn_nodes = [
        ("Node i", Inches(8.2), Inches(2.6), COLOR_CYAN),
        ("Node j", Inches(10.8), Inches(2.3), COLOR_CYAN),
        ("Node u (Detour Insertion)", Inches(9.5), Inches(4.8), COLOR_CYAN)
    ]
    for name, nx, ny, color in gn_nodes:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, nx, ny, Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        
        lbl = slide.shapes.add_textbox(nx - Inches(1.2), ny + Inches(0.35), Inches(2.75), Inches(0.4))
        lbl.text_frame.word_wrap = True
        lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lbl.text_frame.paragraphs[0].text = name
        lbl.text_frame.paragraphs[0].font.size = Pt(11)
        lbl.text_frame.paragraphs[0].font.color.rgb = COLOR_LIGHT_GRAY
        
    # High Prob optimal edge (Thick Cyan line)
    add_route_line(slide, Inches(8.55), Inches(2.77), Inches(10.8), Inches(2.47), COLOR_CYAN, 4.0)
    # Medium Prob edge (Cyan line)
    add_route_line(slide, Inches(10.975), Inches(2.65), Inches(9.675), Inches(4.8), COLOR_CYAN, 2.0)
    # Low Prob edge (Dotted/Thin Red line)
    add_route_line(slide, Inches(8.375), Inches(2.95), Inches(9.5), Inches(4.8), COLOR_DANGER, 1.0)
    
    # Labels with probabilities
    def add_label_box_local(x, y, txt, color):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.8), Inches(0.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_DIM_GRAY
        card.line.width = Pt(1)
        tb = slide.shapes.add_textbox(x, y, Inches(0.8), Inches(0.3))
        tb.text_frame.paragraphs[0].text = txt
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb.text_frame.paragraphs[0].font.size = Pt(10)
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = color
        
    add_label_box_local(Inches(9.4), Inches(2.3), "P=0.92", COLOR_CYAN)
    add_label_box_local(Inches(10.4), Inches(3.6), "P=0.78", COLOR_CYAN)
    add_label_box_local(Inches(8.5), Inches(3.9), "P=0.08", COLOR_DANGER)
    
    set_slide_notes(slide, (
        "Để tăng tốc độ và hướng dẫn tìm kiếm, chúng em tích hợp một mạng thần kinh tích chập đồ thị (GNN) siêu nhẹ.\n\n"
        "Mạng GNN này chạy đúng 1 lần duy nhất lúc bắt đầu để tính toán xác suất liên kết tối ưu giữa mọi cặp khách hàng, "
        "tạo ra một heatmap tĩnh. Chi phí chèn khách hàng mới trong ALNS sẽ được nhân giảm giá dựa trên heatmap này, "
        "giúp ưu tiên các cạnh có độ tin cậy cao.\n\n"
        "Đồng thời, chúng em thiết kế lịch trình lọc cạnh động để bỏ qua sớm các phép biến đổi cục bộ không khả thi, "
        "giúp giảm đáng kể thời gian CPU lãng phí."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 10: POST-PROCESSING SEARCH POLISH CASCADE
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Post-Processing Polish Cascade", "Methodology")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Cascade Optimization Pipeline")
    bullets_s10 = [
        "Depth-3 Ejection Chains: Attempts to empty sparsely-populated routes by displacing customers up to 3 levels deep:",
        "  - Chain: c -> R_i (ejects d) -> d -> R_j (ejects e) -> e -> R_k",
        "Buffered Route Elimination: Triggers multi-route spatial rebalancing if NV = NV_BKS + 1 to force target vehicle termination.",
        "Committed NV Search: Locks the solution to NV - 1 vehicles and executes a strict 1500-iteration SA pathing search with pool-based restarts.",
        "Convergent TD Polish: Executes 2-opt and Or-opt (1, 2, 3) steps along all active routes until distance convergence bounds (Delta c < 10^-9)."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s10, font_size=12.5)
    
    # Right Card
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Ejection Chain Displacement Flow")
    
    # Unassigned node c
    node_c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(3.6), Inches(0.6), Inches(0.6))
    node_c.fill.solid()
    node_c.fill.fore_color.rgb = COLOR_GOLD
    node_c.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(7.5), Inches(3.7), Inches(0.6), Inches(0.4))
    tb.text_frame.paragraphs[0].text = "c"
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    lbl = slide.shapes.add_textbox(Inches(7.0), Inches(4.2), Inches(1.6), Inches(0.4))
    lbl.text_frame.paragraphs[0].text = "Unassigned"
    lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    lbl.text_frame.paragraphs[0].font.size = Pt(10)
    lbl.text_frame.paragraphs[0].font.color.rgb = COLOR_LIGHT_GRAY
    
    # Route R_i card
    add_card(slide, Inches(9.2), Inches(2.2), Inches(1.8), Inches(1.0), border_color=COLOR_CYAN)
    tb_ri = slide.shapes.add_textbox(Inches(9.2), Inches(2.3), Inches(1.8), Inches(0.4))
    tb_ri.text_frame.paragraphs[0].text = "Route R_i"
    tb_ri.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_ri.text_frame.paragraphs[0].font.size = Pt(12)
    tb_ri.text_frame.paragraphs[0].font.bold = True
    tb_ri.text_frame.paragraphs[0].font.color.rgb = COLOR_CYAN
    dot_d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(2.7), Inches(0.3), Inches(0.3))
    dot_d.fill.solid()
    dot_d.fill.fore_color.rgb = COLOR_WHITE
    dot_d.line.fill.background()
    tb_d = slide.shapes.add_textbox(Inches(10.0), Inches(2.7), Inches(0.3), Inches(0.3))
    tb_d.text_frame.paragraphs[0].text = "d"
    tb_d.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_d.text_frame.paragraphs[0].font.size = Pt(9)
    tb_d.text_frame.paragraphs[0].font.bold = True
    tb_d.text_frame.paragraphs[0].font.color.rgb = COLOR_CARD
    
    # Route R_j card
    add_card(slide, Inches(10.3), Inches(4.2), Inches(1.8), Inches(1.0), border_color=COLOR_GOLD)
    tb_rj = slide.shapes.add_textbox(Inches(10.3), Inches(4.3), Inches(1.8), Inches(0.4))
    tb_rj.text_frame.paragraphs[0].text = "Route R_j"
    tb_rj.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_rj.text_frame.paragraphs[0].font.size = Pt(12)
    tb_rj.text_frame.paragraphs[0].font.bold = True
    tb_rj.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD
    dot_e = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.1), Inches(4.7), Inches(0.3), Inches(0.3))
    dot_e.fill.solid()
    dot_e.fill.fore_color.rgb = COLOR_WHITE
    dot_e.line.fill.background()
    tb_e = slide.shapes.add_textbox(Inches(11.1), Inches(4.7), Inches(0.3), Inches(0.3))
    tb_e.text_frame.paragraphs[0].text = "e"
    tb_e.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_e.text_frame.paragraphs[0].font.size = Pt(9)
    tb_e.text_frame.paragraphs[0].font.bold = True
    tb_e.text_frame.paragraphs[0].font.color.rgb = COLOR_CARD

    # Route R_k card (has capacity)
    add_card(slide, Inches(8.2), Inches(5.2), Inches(1.8), Inches(1.0), border_color=COLOR_SUCCESS)
    tb_rk = slide.shapes.add_textbox(Inches(8.2), Inches(5.3), Inches(1.8), Inches(0.4))
    tb_rk.text_frame.paragraphs[0].text = "Route R_k"
    tb_rk.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_rk.text_frame.paragraphs[0].font.size = Pt(12)
    tb_rk.text_frame.paragraphs[0].font.bold = True
    tb_rk.text_frame.paragraphs[0].font.color.rgb = COLOR_SUCCESS
    tb_rk2 = slide.shapes.add_textbox(Inches(8.2), Inches(5.7), Inches(1.8), Inches(0.4))
    tb_rk2.text_frame.paragraphs[0].text = "(Has capacity)"
    tb_rk2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_rk2.text_frame.paragraphs[0].font.size = Pt(10)
    tb_rk2.text_frame.paragraphs[0].font.color.rgb = COLOR_DIM_GRAY

    # Lines showing displacements
    add_route_line(slide, Inches(8.1), Inches(3.6), Inches(9.2), Inches(2.7), COLOR_GOLD, 2)
    add_route_line(slide, Inches(10.1), Inches(3.2), Inches(10.3), Inches(4.2), COLOR_CYAN, 2)
    add_route_line(slide, Inches(11.2), Inches(5.2), Inches(10.0), Inches(5.7), COLOR_GOLD, 2)

    set_slide_notes(slide, (
        "Ở giai đoạn cuối, sau khi thuật toán chính dừng, lời giải tốt nhất tiếp tục đi qua chuỗi 'đánh bóng' (polish cascade).\n\n"
        "Đầu tiên là chuỗi xích đẩy bậc 3 (Depth-3 Ejection Chain): nếu một khách hàng c không thể nhét vào đâu, "
        "ta sẽ chèn c vào tuyến Ri để đẩy d ra, d chèn vào Rj đẩy e ra, e chèn vào Rk. "
        "Chuỗi đẩy sâu 3 cấp này giúp gỡ các nút thắt cực kỳ phức tạp để triệt tiêu chiếc xe cuối cùng.\n\n"
        "Sau đó là các thuật toán loại bỏ xe bằng vùng đệm, tìm kiếm cam kết 1500 vòng lặp, và tối ưu hóa quãng đường cho đến khi hội tụ hoàn toàn."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 11: SYSTEM ARCHITECTURE
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "System Architecture: NAMI Dispatcher", "Technical Stack")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Monorepo Service Layout")
    bullets_s11 = [
        "Frontend (React + Vite): Fast responsive canvas renders spatial routing networks, schedules vehicle slots via interactive Gantt timeline charts.",
        "Backend (FastAPI): Asynchronous API router, streams solver iteration state updates via low-latency WebSocket interfaces.",
        "Heuristic Bindings: Integrates C++ custom local search kernels with SciPy MILP optimizers and Elite state registries.",
        "Structured Data: SQLite database stores benchmark histories, performance logs, and optimized routing profiles."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s11, font_size=13)
    
    # Right Card - Microservice architecture diagram
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Microservice Inter-connects")
    
    # UI Card
    add_card(slide, Inches(7.5), Inches(2.2), Inches(4.6), Inches(0.9), border_color=COLOR_CYAN)
    tb_ui = slide.shapes.add_textbox(Inches(7.5), Inches(2.3), Inches(4.6), Inches(0.7))
    tf = tb_ui.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Vite React Web Portal"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Interactive Canvas | Gantt Timelines"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_CYAN

    # Link lines
    add_route_line(slide, Inches(8.5), Inches(3.1), Inches(8.5), Inches(3.9), COLOR_SUCCESS, 2.5) # WebSockets
    add_route_line(slide, Inches(11.1), Inches(3.1), Inches(11.1), Inches(3.9), COLOR_GOLD, 2)   # HTTP
    
    # WebSocket / HTTP Labels
    tb_lbl1 = slide.shapes.add_textbox(Inches(7.1), Inches(3.3), Inches(1.4), Inches(0.4))
    tb_lbl1.text_frame.paragraphs[0].text = "WS"
    tb_lbl1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_lbl1.text_frame.paragraphs[0].font.size = Pt(10)
    tb_lbl1.text_frame.paragraphs[0].font.bold = True
    tb_lbl1.text_frame.paragraphs[0].font.color.rgb = COLOR_SUCCESS
    
    tb_lbl2 = slide.shapes.add_textbox(Inches(11.1), Inches(3.3), Inches(1.4), Inches(0.4))
    tb_lbl2.text_frame.paragraphs[0].text = "HTTP"
    tb_lbl2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb_lbl2.text_frame.paragraphs[0].font.size = Pt(10)
    tb_lbl2.text_frame.paragraphs[0].font.bold = True
    tb_lbl2.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD

    # API Card
    add_card(slide, Inches(7.5), Inches(3.9), Inches(4.6), Inches(0.9), border_color=COLOR_GOLD)
    tb_api = slide.shapes.add_textbox(Inches(7.5), Inches(4.0), Inches(4.6), Inches(0.7))
    tf = tb_api.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "FastAPI Web Server"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Route Handlers | Socket Broadcasters"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_GOLD
    
    # Link
    add_route_line(slide, Inches(9.8), Inches(4.8), Inches(9.8), Inches(5.4), COLOR_LIGHT_GRAY, 2)
    
    # Core Card
    add_card(slide, Inches(7.5), Inches(5.4), Inches(4.6), Inches(0.9), border_color=COLOR_CYAN)
    tb_core = slide.shapes.add_textbox(Inches(7.5), Inches(5.5), Inches(4.6), Inches(0.7))
    tf = tb_core.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Heuristic Solver Core (Python/C++)"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Hierarchical DDQN | SciPy MILP | ALNS"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_CYAN

    set_slide_notes(slide, (
        "Về mặt kỹ thuật hệ thống, chúng em phát triển nền tảng NAMI Dispatcher dưới dạng cấu trúc monorepo.\n\n"
        "Giao diện frontend viết bằng React và Vite cho tốc độ phản hồi cực nhanh, "
        "vẽ bản đồ tuyến trực quan và hiển thị biểu đồ Gantt thời gian thực.\n\n"
        "Backend sử dụng FastAPI kết nối WebSockets với nhân giải toán Solver Core. "
        "Thiết kế này giúp điều phối viên vừa theo dõi tiến trình giải toán của thuật toán AI từng giây một, "
        "vừa có thể can thiệp thủ công vào các tuyến xe nếu có sự cố ngoài thực tế."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 12: EXPERIMENTAL SETUP
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Experimental Setup & Integrity Protocols", "Evaluation")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Methodological Control")
    bullets_s12 = [
        "Benchmark Target Datasets:",
        "  - Full Solomon 100-customer suite (56 instances).",
        "  - Gehring-Homberger 200 & 400 scaling benchmarks.",
        "Strict Independent Cold-Starts:",
        "  - Every run begins from a freshly cleared directory.",
        "  - Prevents cross-contamination from shared archives (warm-starting from elite routes generated in prior runs), ensuring reproducible evaluations.",
        "Execution Consistency: Worker processes receive identical search boundaries and max iterations via matching CLI parameter overrides to avoid silent defaults."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s12, font_size=12.5)
    
    # Right Card
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Integrity & Hardware Settings")
    
    # Integrity warning box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.2), Inches(4.6), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CARD
    box.line.color.rgb = COLOR_DANGER
    box.line.width = Pt(1.5)
    tb = slide.shapes.add_textbox(Inches(7.7), Inches(2.3), Inches(4.2), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "⚠ Warning on Academic Integrity"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_DANGER
    p.space_after = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = "Several publications report high-performance bounds by running instances sequentially without clearing historical pools. Under independent cold-starts, these benchmarks drop to baseline ranges (e.g. Solomon RC101 drops to 15 vehicles unless proactive seeding is active)."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_LIGHT_GRAY
    
    # Specs box
    add_card(slide, Inches(7.5), Inches(4.7), Inches(4.6), Inches(1.5), border_color=COLOR_DIM_GRAY)
    tb_spec = slide.shapes.add_textbox(Inches(7.7), Inches(4.8), Inches(4.2), Inches(1.3))
    tf = tb_spec.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "Hardware & Tooling Specs"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(6)
    p2 = tf.add_paragraph()
    p2.text = "• CPU: Apple M-Series / Intel Xeon Core-isolated servers.\n• Math Solver: SciPy MILP / Gurobi Engine.\n• Deep Learning: PyTorch (safe-tensors format)."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_LIGHT_GRAY
    
    set_slide_notes(slide, (
        "Để đảm bảo tính trung thực học thuật (Academic Integrity), tụi em áp dụng quy trình kiểm thử cực kỳ nghiêm ngặt "
        "gọi là 'Independent Cold-Starts'.\n\n"
        "Nhiều bài báo trước đây thường chạy liên tục các thực thể, khiến lời giải tốt của thực thể trước nằm lại trong bộ nhớ đệm "
        "và vô tình 'mồi' (warm-start) cho thực thể sau, dẫn đến kết quả sai lệch và không thể tái lập. "
        "Ở đây, mỗi lượt giải của tụi em đều được cô lập hoàn toàn từ đầu, thư mục bộ đệm bị xóa sạch trước khi chạy.\n\n"
        "Thực nghiệm được tiến hành trên toàn bộ 56 thực thể Solomon 100 khách hàng, và các quy mô lớn 200, 400 khách hàng của Gehring-Homberger."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 13: EXPERIMENTAL RESULTS - VEHICLE COUNT
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Solomon Benchmark: Vehicle Minimization", "Results")
    
    # Table Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0), "Average Fleet Size Comparison (Mean NV)")
    
    # Table
    rows = 7
    cols = 6
    table_shape = slide.shapes.add_table(rows, cols, Inches(1.1), Inches(2.2), Inches(11.1), Inches(3.8))
    table = table_shape.table
    
    # Headers
    headers = ["Instance Class", "BKS Value", "Google OR-Tools", "ALNS-Base Heuristic", "Hybrid DDQN-ALNS (Ours)", "Improvement vs OR-Tools"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    data = [
        ["Solomon C1", "10.00", "10.00", "10.00", "10.00", "Matched"],
        ["Solomon R1", "11.92", "12.08", "12.00", "11.96", "-0.99%"],
        ["Solomon RC1", "11.50", "11.88", "11.63", "11.46*", "-3.53%"],
        ["Solomon C2", "3.00", "3.00", "3.00", "3.00", "Matched"],
        ["Solomon R2", "2.73", "3.18", "2.82", "2.73", "-14.15%"],
        ["Solomon RC2", "3.25", "6.25 †", "3.50", "3.47", "-44.48%"]
    ]
    
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_LIGHT_GRAY
            p.alignment = PP_ALIGN.CENTER
            
            # Highlight ours
            if c_idx == 4:
                p.font.bold = True
                p.font.color.rgb = COLOR_SUCCESS
            if c_idx == 5 and "-" in val:
                p.font.bold = True
                p.font.color.rgb = COLOR_CYAN
            if "†" in val:
                p.font.bold = True
                p.font.color.rgb = COLOR_DANGER
                
    # Footnote Box
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(6.1), Inches(11.1), Inches(0.4))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = "* Note: Negative gaps represent bounds beating or matching BKS. † Marker highlights vehicle inflation in standard solvers on tight time bounds."
    p.font.name = FONT_BODY
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_DIM_GRAY
    
    set_slide_notes(slide, (
        "Đây là kết quả thực nghiệm về số xe — mục tiêu quan trọng nhất.\n\n"
        "Trên nhóm thực thể RC1, thuật toán lai DDQN-ALNS của chúng em đạt Gap% âm so với BKS (-0.358%), "
        "chứng tỏ chúng em tìm được số xe tối thiểu toàn cục (ví dụ như mức 14 xe của RC101) trong khi ALNS thông thường bị kẹt ở 15 xe.\n\n"
        "Đặc biệt trên nhóm RC2, thuật toán của tụi em duy trì số xe trung bình cực tốt là 3.475 xe, "
        "trong khi bộ giải Google OR-Tools bị hiện tượng lạm phát xe nghiêm trọng khi phải dùng trung bình tới 6.25 xe "
        "(tức là tốn gần gấp đôi số xe) để đổi lấy quãng đường ngắn hơn."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 14: EXPERIMENTAL RESULTS - DISTANCE & FAIR COMPARISON
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Strict Fair Intersection Routing Quality", "Results")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "The Fallacy of Raw Travel Distance Gaps")
    bullets_s14 = [
        "Distorted Comparisons: Travel Distance (TD) comparisons are mathematically flawed if vehicle counts are not matched. Extra vehicle capacity relaxes spatial routing limits, artificially lowering travel distance.",
        "Strict Fair Intersection Protocol (N = 39):",
        "  - Isolates the subset of 39 Solomon instances where both solvers successfully matched BKS vehicle counts.",
        "  - Guarantees a mathematically direct comparison of routing efficiency under identical time constraints.",
        "  - Hybrid DDQN-ALNS outperforms ALNS-Base across all families."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s14, font_size=13.5)
    
    # Right Card
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "TD Improvement Gaps (Fair Intersects)")
    
    # Box 1
    add_card(slide, Inches(7.5), Inches(2.3), Inches(4.6), Inches(1.8), border_color=COLOR_SUCCESS)
    tb1 = slide.shapes.add_textbox(Inches(7.7), Inches(2.5), Inches(4.2), Inches(1.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "RC1 Family Gaps"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf1.add_paragraph()
    p2.text = "-1.75%"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_SUCCESS
    p3 = tf1.add_paragraph()
    p3.text = "Travel Distance reduction compared directly to ALNS-Base."
    p3.font.name = FONT_BODY
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_LIGHT_GRAY
    
    # Box 2
    add_card(slide, Inches(7.5), Inches(4.4), Inches(4.6), Inches(1.8), border_color=COLOR_SUCCESS)
    tb2 = slide.shapes.add_textbox(Inches(7.7), Inches(4.6), Inches(4.2), Inches(1.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "R1 Family Gaps"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p2 = tf2.add_paragraph()
    p2.text = "-4.07%"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_SUCCESS
    p3 = tf2.add_paragraph()
    p3.text = "Travel Distance reduction compared directly to ALNS-Base."
    p3.font.name = FONT_BODY
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_LIGHT_GRAY
    
    set_slide_notes(slide, (
        "Ở slide này, em xin nhấn mạnh một đóng góp về phương pháp luận đánh giá.\n\n"
        "Khi so sánh quãng đường di chuyển (TD), việc so sánh thô là không sòng phẳng nếu số xe không bằng nhau. "
        "Một thuật toán dùng 15 xe chắc chắn sẽ dễ tìm đường ngắn hơn thuật toán dùng 14 xe vì áp lực ràng buộc thời gian được san sẻ.\n\n"
        "Vì vậy, tụi em áp dụng bộ lọc 'Strict Fair Intersection' (N=39) — tức là chỉ so sánh quãng đường "
        "trên các thực thể mà tất cả các bộ giải đều đạt cùng số xe tối thiểu. Kết quả cho thấy thuật toán lai của chúng em "
        "giảm được 1.75% quãng đường trên nhóm RC1 và tới 4.07% trên nhóm R1 so với ALNS truyền thống."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 15: GEHRING--HOMBERGER SCALABILITY & LIMITATIONS
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Homberger Scalability & Computational Overhead", "Evaluation")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Large-Scale Routing Profiles")
    bullets_s15 = [
        "200-Customer Scale (Robust Performance):",
        "  - Both solvers converge to identical vehicle limits.",
        "  - Hybrid-DDQN limits vehicle degradation rates to 0% - 20% compared to 30% - 70% for ALNS-Base.",
        "400-Customer Scale (Suboptimal Degradation):",
        "  - Neither solver approaches the BKS limits (e.g. BKS NV=4 on R2_4_1, solvers land at 8.10 - 8.80).",
        "  - However, Hybrid DDQN-ALNS yields a statistically significant edge of 0.70-0.80 vehicles (Wilcoxon p=0.0156 on R2_4_1)."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s15, font_size=13.0)
    
    # Right Card
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Computational Performance Limits")
    
    tb_limit = slide.shapes.add_textbox(Inches(7.4), Inches(2.2), Inches(4.8), Inches(4.0))
    tf = tb_limit.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "Computational Overhead"
    p.font.name = FONT_TITLE
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = "The hybrid solver is a quality-maximizing strategy that trades CPU cycles for routing accuracy. It is not an execution speed accelerator.\n\nRun Overhead Metrics:"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_LIGHT_GRAY
    p2.space_after = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "• Solomon-100: 1.5x - 4.0x slower runtime.\n• Homberger-200: 1.3x - 4.6x slower runtime.\n• Homberger-400: Up to 100x slower on pathological configurations due to large dense MILP matrices (specifically R2_4_1)."
    p3.font.name = FONT_BODY
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_GOLD
    
    set_slide_notes(slide, (
        "Về khả năng mở rộng trên tập Gehring-Homberger quy mô lớn:\n\n"
        "Ở quy mô 200 khách hàng, thuật toán lai thể hiện độ ổn định vượt trội khi tỷ lệ lệch số xe chỉ từ 0% đến 20%, "
        "so với mức 30% đến 70% của ALNS truyền thống.\n\n"
        "Ở quy mô 400 khách hàng, dù cả hai thuật toán đều có sự suy giảm so với BKS lý thuyết, "
        "mô hình của chúng em vẫn giữ được vị thế dẫn trước rõ rệt với việc tiết kiệm trung bình từ 0.7 đến 0.8 xe "
        "(có ý nghĩa thống kê với kiểm định Wilcoxon p=0.0156).\n\n"
        "Đổi lại, thuật toán lai chạy chậm hơn từ 1.3 đến 4.6 lần do chi phí giải các bài toán MILP phân hoạch tập hợp lớn."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 16: SYSTEM DEMONSTRATION
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "NAMI Dispatcher Telemetry Demonstration", "System Showcase")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Telemetry & Operator Controls")
    bullets_s16 = [
        "Live Telemetry: Iteration-by-iteration progress metrics (NV, TD, convergence gap) are streamed from the FastAPI solver core to React web interfaces via WebSockets.",
        "Gantt Schedule Editor: Visualizes arrival times, service durations, and slack periods. Allows dispatchers to drag-and-drop stops to edit schedules manually.",
        "Feasibility Monitors: Dynamically flags capacity and time window conflicts in real time as the dispatcher edits routes."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s16, font_size=13.5)
    
    # Right Card - UI Mockup
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Interactive Interface Mockup")
    
    # Draw App Shell
    add_card(slide, Inches(7.4), Inches(2.2), Inches(4.8), Inches(4.0), border_color=COLOR_DIM_GRAY)
    
    # Inner elements
    # Map box
    add_card(slide, Inches(7.6), Inches(2.4), Inches(2.1), Inches(2.0), border_color=COLOR_CYAN)
    tb = slide.shapes.add_textbox(Inches(7.6), Inches(2.5), Inches(2.1), Inches(0.4))
    tb.text_frame.paragraphs[0].text = "Spatial Map"
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb.text_frame.paragraphs[0].font.size = Pt(11)
    tb.text_frame.paragraphs[0].font.bold = True
    
    # Telemetry Box
    add_card(slide, Inches(9.9), Inches(2.4), Inches(2.1), Inches(2.0), border_color=COLOR_GOLD)
    tb2 = slide.shapes.add_textbox(Inches(9.9), Inches(2.5), Inches(2.1), Inches(0.4))
    tb2.text_frame.paragraphs[0].text = "Live Telemetry"
    tb2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb2.text_frame.paragraphs[0].font.size = Pt(11)
    tb2.text_frame.paragraphs[0].font.bold = True
    
    # Gantt Box
    add_card(slide, Inches(7.6), Inches(4.6), Inches(4.4), Inches(1.2), border_color=COLOR_SUCCESS)
    tb3 = slide.shapes.add_textbox(Inches(7.6), Inches(4.7), Inches(4.4), Inches(0.4))
    tb3.text_frame.paragraphs[0].text = "Gantt Timeline Editor (Schedule Slacks)"
    tb3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb3.text_frame.paragraphs[0].font.size = Pt(11)
    tb3.text_frame.paragraphs[0].font.bold = True

    set_slide_notes(slide, (
        "Đây là giao diện của hệ thống trực quan hóa NAMI Dispatcher do tụi em xây dựng.\n\n"
        "Hệ thống cho phép điều phối viên tải dữ liệu khách hàng lên, bấm nút giải toán và theo dõi "
        "tiến trình chạy của thuật toán AI dưới dạng đồ thị telemetry trực quan.\n\n"
        "Bản đồ sẽ vẽ chi tiết lộ trình của từng xe với màu sắc khác nhau, và biểu đồ Gantt bên dưới giúp kiểm tra "
        "độ trễ (slack time) của từng điểm giao hàng. Công cụ này giúp biến thuật toán AI dạng hộp đen thành một hệ thống minh bạch, "
        "có thể tương tác và ứng dụng được ngay vào doanh nghiệp."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 17: CONCLUSION
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Conclusion & Future Directions", "Summary")
    
    # Left Card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "Key Contributions")
    bullets_s17 = [
        "State-Conditioned Search: Validated the utility of hierarchical MDP modeling to direct local search operators.",
        "Resolved Vehicle Limits: Successfully bypassed local vehicle minimization bounds on Solomon benchmarks (solving hard instances like RC101 to global minimums).",
        "Operational Dispatch Tooling: Created NAMI Dispatcher, bridging the gap between deep mathematical optimization and dispatcher software interfaces."
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.4), Inches(4.0), bullets_s17, font_size=13.5)
    
    # Right Card
    add_card(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(5.0), "Future Enhancements")
    bullets_s17_r = [
        "Multi-GPU Parallel Workloads: Scaling training rollouts across multi-GPU setups to decrease offline policy training times.",
        "Heterogeneous Fleet Support: Generalizing constraint layouts to handle vehicles with varying capacities and costs.",
        "Dynamic Real-time Routing: Adjusting optimization weights dynamically to adapt to traffic accidents or delivery delays."
    ]
    add_bullet_list(slide, Inches(7.4), Inches(2.3), Inches(4.8), Inches(4.0), bullets_s17_r, font_size=13.5)
    
    set_slide_notes(slide, (
        "Tóm lại, nghiên cứu của tụi em đã:\n"
        "1. Chứng minh tính hiệu quả của mô hình học máy phân cấp trong điều khiển tìm kiếm cục bộ;\n"
        "2. Vượt qua trần xe tối thiểu trên các bộ dữ liệu chuẩn Solomon;\n"
        "3. Và hiện thực hóa giải pháp bằng một phần mềm trực quan hoàn chỉnh.\n\n"
        "Trong tương lai, nhóm sẽ tập trung tối ưu hóa thời gian tính toán trên quy mô lớn bằng cách song song hóa luồng chạy trên GPU, "
        "và mở rộng bài toán cho đội xe không đồng nhất (nhiều tải trọng khác nhau).\n\n"
        "Em xin chân thành cảm ơn thầy cô đã lắng nghe và rất mong nhận được câu hỏi góp ý từ Hội đồng."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 18: DEFENSE STRATEGIES (Q&A PANEL)
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide)
    add_header(slide, "Defense Strategy Panel (Q&A Guidelines)", "Q&A Preparation")
    
    # Full size card
    add_card(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0), "Pre-empting Academic Review Committee Objections")
    
    tb_qa = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(11.1), Inches(4.2))
    tf = tb_qa.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Q1
    p = tf.paragraphs[0]
    p.text = "Q1: Why is the hybrid solver execution runtime slower than standard ALNS?"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(4)
    
    p_a1 = tf.add_paragraph()
    p_a1.text = "A1: It represents a deliberate, cost-minimizing trade-off. For logistics companies, saving a single vehicle reduces fixed monthly overheads (capital depreciation, licensing, driver salaries) by thousands of dollars. Investing an extra 2-3 minutes of server computation before vehicles dispatch yields substantial cost reductions compared to starting immediately with an extra truck."
    p_a1.font.name = FONT_BODY
    p_a1.font.size = Pt(11)
    p_a1.font.color.rgb = COLOR_LIGHT_GRAY
    p_a1.space_after = Pt(12)
    
    # Q2
    p2 = tf.add_paragraph()
    p2.text = "Q2: Why do you report GNN performance checks on only 3 instances (C101, R101, RC101)?"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_CYAN
    p2.space_after = Pt(4)
    
    p_a2 = tf.add_paragraph()
    p_a2.text = "A2: The GNN evaluations in Section IV-F serve as a representative case study to verify that neural heatmaps properly adjust to distinct spatial configurations (clustered, random, mixed). Due to deep execution costs, we focused on trace profiling. Complete results showing average metrics across all 56 instances are reported in our main thesis paper."
    p_a2.font.name = FONT_BODY
    p_a2.font.size = Pt(11)
    p_a2.font.color.rgb = COLOR_LIGHT_GRAY
    p_a2.space_after = Pt(12)
    
    # Q3
    p3 = tf.add_paragraph()
    p3.text = "Q3: How do we know these solver improvements are mathematically significant and not random?"
    p3.font.name = FONT_TITLE
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_SUCCESS
    p3.space_after = Pt(4)
    
    p_a3 = tf.add_paragraph()
    p_a3.text = "A3: All comparative distance assertions are validated using the non-parametric Wilcoxon Signed-Rank Test at a strict significance boundary (alpha = 0.05). The improvements are accompanied by explicit statistical validation (e.g. p = 0.0078 on instance c2_4_1), proving the hybrid solver benefits are mathematically significant."
    p_a3.font.name = FONT_BODY
    p_a3.font.size = Pt(11)
    p_a3.font.color.rgb = COLOR_LIGHT_GRAY
    
    set_slide_notes(slide, (
        "Dưới đây là cẩm nang trả lời các câu hỏi phản biện cốt lõi từ Hội đồng khoa học:\n\n"
        "1. Về thời gian tính toán: Nhấn mạnh đây là sự đánh đổi có lợi giữa tài nguyên máy tính giá rẻ "
        "và chi phí vận hành xe thực tế hàng ngàn USD của doanh nghiệp.\n\n"
        "2. Về việc đánh giá GNN trên 3 thực thể: Giải thích rõ đây là Case Study để kiểm định tính thích ứng cấu trúc không gian "
        "(clustered, random, mixed), và khẳng định số liệu quét đầy đủ đã nằm trong báo cáo chi tiết.\n\n"
        "3. Về độ tin cậy khoa học: Nhắc đến việc tất cả các kết quả so sánh đều đi kèm kiểm định Wilcoxon Signed-Rank "
        "với p-value rõ ràng chứng minh độ vượt trội mang ý nghĩa thống kê."
    ))

    # ═══════════════════════════════════════════════════════════════════
    # SAVE PRESENTATION
    # ═══════════════════════════════════════════════════════════════════
    output_path = os.path.join("docs", "slides.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_slides()
