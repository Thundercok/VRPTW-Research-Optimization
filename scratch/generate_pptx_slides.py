import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    BG_DARK = RGBColor(11, 17, 34)       # #0b1122
    CARD_BG = RGBColor(18, 26, 47)       # #121a2f
    PRIMARY = RGBColor(0, 180, 216)      # Cyan
    ACCENT = RGBColor(245, 158, 11)     # Amber
    PURPLE = RGBColor(114, 9, 183)      # Purple
    GREEN = RGBColor(16, 185, 129)      # Emerald
    TEXT_WHITE = RGBColor(248, 250, 252) # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)# #94a3b8
    
    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg
        
    def add_header(slide, title, category="NÂNG CAO NĂNG LỰC TỐI ƯU HÓA VRPTW"):
        # Category badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, title="", bg_color=CARD_BG, border_color=PRIMARY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        
        if title:
            tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = PRIMARY
        return card

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)
    add_card(slide1, 0.8, 1.2, 11.733, 5.1, bg_color=CARD_BG, border_color=PRIMARY)
    
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.933), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "BÁO CÁO NGHIÊN CỨU KHOA HỌC CẤP TRƯỜNG"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT
    
    p1 = tf.add_paragraph()
    p1.text = "Hệ Thống Tối Ưu Lộ Trình Xe Có Khung Thời Gian (VRPTW)\nBằng Thuật Toán Lai Deep Reinforcement Learning (DDQN) & ALNS"
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_before = Pt(15)
    
    p2 = tf.add_paragraph()
    p2.text = "A Hybrid Deep Q-Network & Adaptive Large Neighborhood Search Framework with MILP Recombination"
    p2.font.size = Pt(15)
    p2.font.italic = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "• Thời lượng báo cáo: 20 Phút  |  • Quy mô thực nghiệm: Solomon 100 & Gehring-Homberger 200-600\n• Đóng góp: Tối ưu số xe (NV), Giảm quãng đường (TD), Kiểm định thống kê Wilcoxon p < 0.002"
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(30)

    # SLIDE 2: Agenda & Executive Summary
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_header(slide2, "Nội Dung Báo Cáo (Agenda - 20 Phút)", "TỔNG QUAN BÀI THUYẾT TRÌNH")
    
    items = [
        ("1. Đặt Vấn Đề & Mục Tiêu", "Bài toán VRPTW, mục tiêu thứ bậc Lexicographic (NV > TD) & Thách thức scale."),
        ("2. Tổng Quan Nghiên Cứu", "So sánh ALNS cổ điển, Google OR-Tools & Khoảng trống nghiên cứu (Research Gap)."),
        ("3. Phương Pháp Đề Xuất (Core)", "Khung thuật toán lai Hybrid-DDQN: Regime Switching, Operator RL, GNN Edge & MILP."),
        ("4. Kiến Trúc Hệ Thống & Web App", "Full-stack Monorepo, FastAPI, PyTorch, HiGHS MILP, Firebase & UI Visualizer Dashboard."),
        ("5. Đánh Giá Thực Nghiệm", "So sánh độc lập (Cold-start), Kiểm định Wilcoxon (p < 0.002), Phân tích lạm phát xe."),
        ("6. Kết Luận & Đóng Góp", "Tóm tắt đóng góp chính & Định hướng phát triển ứng dụng thực tế.")
    ]
    
    for idx, (title, desc) in enumerate(items):
        r = idx // 3
        c = idx % 3
        left = 0.8 + c * 3.97
        top = 1.6 + r * 2.6
        add_card(slide2, left, top, 3.75, 2.3, title=title, border_color=PRIMARY if idx==2 else CARD_BG)
        
        tb = slide2.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.7), Inches(3.35), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED

    # SLIDE 3: Problem Definition
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_header(slide3, "1. Đặt Vấn Đề: Bài Toán VRPTW & Mục Tiêu Lexicographic", "PHẦN 1: GIỚI THIỆU & ĐỘNG LỰC")
    
    add_card(slide3, 0.8, 1.6, 5.7, 5.1, title="Mô Hình Bài Toán VRPTW", border_color=PRIMARY)
    tb1 = slide3.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.3), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "• Định nghĩa: Phục vụ N khách hàng với vị trí, nhu cầu d_i, thời gian phục vụ s_i và khung thời gian [e_i, l_i]."
    tf1.paragraphs[0].font.size = Pt(14)
    tf1.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    p = tf1.add_paragraph()
    p.text = "• Đồ thị G = (V, A): V = {0, 1, ..., N}, 0 là kho trung tâm (depot)."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(10)
    
    p = tf1.add_paragraph()
    p.text = "• Ràng buộc nghiêm ngặt:\n  - Khung thời gian: e_i <= t_i <= l_i\n  - Sức chứa xe: sum(d_i) <= Capacity Q\n  - Ràng buộc quay về kho: t_last + s_last + dist(last, 0) <= l_0"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(10)

    add_card(slide3, 6.8, 1.6, 5.7, 5.1, title="Hàm Mục Tiêu Lexicographic (NV >> TD)", border_color=ACCENT)
    tb2 = slide3.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.3), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "Hàm mục tiêu ưu tiên bậc nhất trong Thực Tế Logistics:"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = ACCENT
    
    p = tf2.add_paragraph()
    p.text = "Minimizing:  min ( NV , TD )"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_before = Pt(10)
    
    p = tf2.add_paragraph()
    p.text = "1. Ưu tiên hàng đầu (Primary): Giảm số lượng xe deployed (NV).\n   -> Chi phí cố định đầu xe (khấu hao, tài xế, bảo trì) áp đảo chi phí nhiên liệu di chuyển.\n2. Ưu tiên thứ hai (Secondary): Tối thiểu hóa tổng quãng đường di chuyển (TD) khi NV đã cố định."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(15)

    # SLIDE 4: Proposed Method Architecture
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_header(slide4, "3. Phương Pháp Đề Xuất: Kiến Trúc Lai Hybrid DDQN-ALNS", "PHẦN 3: ĐÓNG GÓP THUẬT TOÁN")
    
    add_card(slide4, 0.8, 1.6, 11.733, 5.1, title="Khung Thuật Toán Điều Khiển Phân Cấp (Hierarchical DRL Control Loop)", border_color=PRIMARY)
    tb = slide4.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.133), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "1. Plateau Controller (DDQN Trạng Thái Vùng Tìm Kiếm):\n   - Quyết định chuyển đổi 7 Chế Độ (Regime Switching): Intensify, Diversify, Route-Reduce, Pool-Recombine, TW-Rescue...\n   - Giúp tìm kiếm bứt phá khỏi vùng cực trị địa phương."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    
    p = tf.add_paragraph()
    p.text = "2. Operator Controller (DDQN Lựa Chọn Toán Tử):\n   - Lựa chọn thông minh cặp toán tử trong không gian 65 hành động (13 Destroy x 5 Repair).\n   - Kết hợp mô hình Learned Acceptance Criterion (LAC) thay thế Metropolis Simulated Annealing."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(15)

    p = tf.add_paragraph()
    p.text = "3. GNN Edge Heatmap Guidance & Generalized Ejection Chains:\n   - GNN dự đoán xác suất cạnh tối ưu để định hướng sửa lộ trình.\n   - Ejection Chains độ sâu N (generalized depth 4-6) giúp triệt tiêu các route khó."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(15)

    # SLIDE 5: Benchmark Table
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_header(slide5, "5. Kết Quả Thực Nghiệm Trên Bộ Dữ Liệu Solomon 100", "PHẦN 5: ĐÁNH GIÁ THỰC NGHIỆM")
    
    add_card(slide5, 0.8, 1.6, 11.733, 5.1, title="Bảng Tổng Hợp Hiệu Năng Solomon 100 (56 Thực Thể)", border_color=PRIMARY)
    
    rows, cols = 7, 6
    left, top, width, height = Inches(1.1), Inches(2.3), Inches(11.133), Inches(4.2)
    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Lớp Dữ Liệu", "Thuật Toán", "Mean NV", "Mean TD", "Gap% vs BKS", "Thời Gian (s)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = BG_DARK

    data = [
        ["C1 (17 bài)", "Hybrid-DDQN", "10.00", "828.38", "0.00%", "95.2s"],
        ["C2 (8 bài)", "Hybrid-DDQN", "3.00", "589.97", "+0.02%", "62.4s"],
        ["R1 (12 bài)", "Hybrid-DDQN", "12.26", "1206.81", "-0.37%", "117.6s"],
        ["R2 (11 bài)", "Hybrid-DDQN", "2.86", "950.61", "-0.08%", "125.2s"],
        ["RC1 (8 bài)", "Hybrid-DDQN", "11.98", "1358.82", "-1.50%", "113.5s"],
        ["RC2 (8 bài)", "Hybrid-DDQN", "3.25", "1130.02", "+3.18%", "78.3s"],
    ]
    
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_WHITE

    # SLIDE 6: Statistical Significance
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_header(slide6, "5. Kiểm Định Thống Kê Wilcoxon Signed-Rank Test", "PHẦN 5: ĐÁNH GIÁ THỰC NGHIỆM")
    
    add_card(slide6, 0.8, 1.6, 5.7, 5.1, title="Kiểm Định Số Lượng Xe (NV)", border_color=PRIMARY)
    tb1 = slide6.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.3), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "• Mẫu kiểm định N = 56 bài Solomon."
    tf1.paragraphs[0].font.size = Pt(13)
    tf1.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    p = tf1.add_paragraph()
    p.text = "Wilcoxon p-value = 0.00161"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_before = Pt(15)
    
    p = tf1.add_paragraph()
    p.text = "-> Ý nghĩa thống kê vượt trội (p < 0.01):\n  Hybrid-DDQN giảm số lượng xe đáng kể so với ALNS-Base."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(10)

    add_card(slide6, 6.8, 1.6, 5.7, 5.1, title="Kiểm Định Quãng Đường (Matched-NV TD)", border_color=ACCENT)
    tb2 = slide6.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.3), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "• Mẫu kiểm định N = 45 bài (Khớp chính xác số xe)."
    tf2.paragraphs[0].font.size = Pt(13)
    tf2.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    p = tf2.add_paragraph()
    p.text = "Wilcoxon p-value = 0.000415"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_before = Pt(15)
    
    p = tf2.add_paragraph()
    p.text = "-> Ý nghĩa thống kê cực kỳ cao (p < 0.001):\n  Khi cùng số xe, Hybrid-DDQN tối ưu quãng đường vượt trội so với ALNS truyền thống."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(10)

    # Save presentation
    output_path = "docs/slides_v2.pptx"
    prs.save(output_path)
    print(f"Presentation deck successfully saved to: {output_path}")

if __name__ == "__main__":
    create_deck()
